from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import logging
import os
import re  # Add this import for regex
from werkzeug.utils import secure_filename
import PyPDF2
import traceback
from pathlib import Path
import torch
from transformers import GPT2LMHeadModel, GPT2Tokenizer, AutoModelForCausalLM, AutoTokenizer
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from collections import defaultdict
import random
from docx import Document  # For .docx files
from pptx import Presentation  # For .pptx files
from datetime import datetime, timedelta
from peft import (
    LoraConfig,
    TaskType,
    get_peft_model,
    PeftModel,
    PeftConfig
)
from pymongo import MongoClient
from bson.objectid import ObjectId
import os.path
from reportlab.lib.units import inch
import time

# Configure logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger('werkzeug')
logger.setLevel(logging.DEBUG)

# Create Flask app
app = Flask(__name__)

# Configure CORS
CORS(app, resources={
    r"/*": {
        "origins": ["http://localhost:3000", "http://127.0.0.1:3000"],
        "methods": ["GET", "POST", "OPTIONS"],
        "allow_headers": ["Content-Type", "Authorization"]
    }
})

# Configuration
app.config.update({
    'MAX_CONTENT_LENGTH': 50 * 1024 * 1024,  # Increase to 50MB
    'UPLOAD_FOLDER': 'uploads',
    'SUMMARIES_DIR': 'summaries',
    'ALLOWED_EXTENSIONS': {'txt', 'pdf', 'docx', 'pptx'},
    'MODEL_NAME': 'gpt2',
    'PAPERS_DIR': os.path.join(os.path.dirname(os.path.abspath(__file__)), 'papers')
})

# Ensure directories exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['SUMMARIES_DIR'], exist_ok=True)

# Global variables for models
base_model = None
peft_model = None
tokenizer = None
current_context_hash = None

# Update ML domain knowledge with more comprehensive keywords and variations
ML_CONCEPTS = {
    'neural_networks': {
        'keywords': [
            'neural network', 'neural nets', 'deep learning', 'deep neural network',
            'CNN', 'RNN', 'LSTM', 'GRU', 'transformer', 'attention mechanism',
            'backpropagation', 'forward propagation', 'weights', 'biases',
            'activation function', 'relu', 'sigmoid', 'tanh', 'softmax',
            'hidden layer', 'output layer', 'input layer', 'neuron', 'perceptron',
            'deep network', 'feedforward', 'convolutional', 'recurrent'
        ],
        'relationships': ['training', 'inference', 'learning', 'optimization'],
        'definitions': {
            'neural network': 'A computational model inspired by biological neural networks that learns patterns from data',
            'deep learning': 'A subset of machine learning using neural networks with multiple layers',
            'CNN': 'Convolutional Neural Network - specialized for processing grid-like data such as images',
            'RNN': 'Recurrent Neural Network - designed to work with sequential data',
            'backpropagation': 'Algorithm for training neural networks by computing gradients of the loss function'
        }
    },
    'convolutional_networks': {
        'keywords': [
            'CNN', 'ConvNet', 'convolutional neural network', 'convolution layer',
            'pooling layer', 'max pooling', 'average pooling', 'feature map',
            'kernel', 'filter', 'stride', 'padding', 'channels', 'receptive field',
            'feature extraction', 'image recognition', 'object detection',
            'VGG', 'ResNet', 'Inception', 'transfer learning', 'fine-tuning',
            'batch normalization', 'dropout', 'data augmentation'
        ],
        'relationships': ['image processing', 'feature extraction', 'classification'],
        'definitions': {
            'CNN': 'Neural network architecture specialized for processing grid-like data, particularly effective for image analysis',
            'convolution layer': 'Layer that applies filters to input data to detect features and patterns',
            'pooling layer': 'Layer that reduces spatial dimensions while preserving important features',
            'feature map': 'Output of a convolution operation showing detected features in the input'
        }
    },
    'transformers': {
        'keywords': [
            'transformer', 'attention mechanism', 'self-attention', 'multi-head attention',
            'encoder', 'decoder', 'positional encoding', 'BERT', 'GPT', 'T5',
            'language model', 'sequence-to-sequence', 'token embedding',
            'masked attention', 'cross attention', 'transformer block',
            'layer normalization', 'feed forward network', 'residual connection',
            'pre-training', 'fine-tuning', 'transfer learning'
        ],
        'relationships': ['sequence modeling', 'language processing', 'attention computation'],
        'definitions': {
            'transformer': 'Neural network architecture that uses self-attention mechanisms to process sequential data',
            'attention mechanism': 'Technique allowing models to focus on relevant parts of input data',
            'self-attention': 'Method for computing relationships between different positions in a sequence',
            'BERT': 'Bidirectional Encoder Representations from Transformers - a transformer-based language model'
        }
    },
    'generative_models': {
        'keywords': [
            'GAN', 'generative adversarial network', 'generator', 'discriminator',
            'adversarial training', 'latent space', 'style transfer', 'image synthesis',
            'DCGAN', 'WGAN', 'StyleGAN', 'CycleGAN', 'conditional GAN',
            'mode collapse', 'nash equilibrium', 'minimax game',
            'progressive growing', 'adversarial loss', 'reconstruction loss',
            'variational autoencoder', 'VAE', 'diffusion models'
        ],
        'relationships': ['generation', 'synthesis', 'adversarial training'],
        'definitions': {
            'GAN': 'Framework where two networks compete to generate realistic data',
            'generator': 'Network that creates synthetic data samples',
            'discriminator': 'Network that distinguishes between real and synthetic data',
            'latent space': 'Compressed representation space used for generation'
        }
    },
    'ensemble_learning': {
        'keywords': [
            'ensemble', 'bagging', 'boosting', 'random forest', 'voting',
            'stacking', 'weak learner', 'strong learner', 'bootstrap',
            'AdaBoost', 'XGBoost', 'LightGBM', 'CatBoost', 'gradient boosting',
            'majority voting', 'weighted voting', 'model averaging',
            'diversity', 'bias-variance trade-off', 'meta-learning',
            'cross-validation', 'out-of-bag error', 'feature importance'
        ],
        'relationships': ['combination', 'aggregation', 'voting', 'boosting'],
        'definitions': {
            'ensemble learning': 'Technique combining multiple models to improve prediction accuracy',
            'bagging': 'Bootstrap aggregating - training models on random subsets of data',
            'boosting': 'Sequential training of models focusing on hard examples',
            'random forest': 'Ensemble of decision trees using bagging and random feature selection'
        }
    },
    'decision_trees': {
        'keywords': [
            'decision tree', 'root node', 'leaf node', 'splitting criterion',
            'information gain', 'gini index', 'entropy', 'pruning',
            'CART', 'ID3', 'C4.5', 'binary split', 'multi-way split',
            'tree depth', 'minimum samples split', 'maximum features',
            'feature importance', 'decision boundary', 'regression tree',
            'classification tree', 'recursive partitioning'
        ],
        'relationships': ['splitting', 'classification', 'regression'],
        'definitions': {
            'decision tree': 'Tree-like model making decisions based on feature conditions',
            'information gain': 'Measure of reduction in entropy after a split',
            'pruning': 'Process of removing branches to prevent overfitting',
            'CART': 'Classification and Regression Trees algorithm'
        }
    },
    'deep_learning': {
        'keywords': [
            'deep learning', 'neural architecture', 'deep neural network',
            'representation learning', 'feature hierarchy', 'end-to-end learning',
            'transfer learning', 'pre-training', 'fine-tuning', 'deep belief network',
            'autoencoder', 'regularization', 'optimization', 'gradient descent',
            'backpropagation', 'activation function', 'loss function',
            'batch normalization', 'dropout', 'residual connection',
            'skip connection', 'dense layer', 'embedding layer'
        ],
        'relationships': ['hierarchical learning', 'feature extraction', 'optimization'],
        'definitions': {
            'deep learning': 'Subset of machine learning using multiple layers to learn representations',
            'representation learning': 'Learning meaningful features automatically from raw data',
            'transfer learning': 'Reusing knowledge learned from one task for another task',
            'end-to-end learning': 'Learning directly from raw input to desired output without manual feature engineering'
        }
    },
    'optimization_techniques': {
        'keywords': [
            'gradient descent', 'stochastic gradient descent', 'mini-batch',
            'momentum', 'Adam', 'RMSprop', 'AdaGrad', 'learning rate',
            'learning rate schedule', 'weight decay', 'L1 regularization',
            'L2 regularization', 'dropout', 'batch normalization',
            'early stopping', 'model checkpoint', 'gradient clipping',
            'optimizer', 'loss function', 'convergence', 'local minima',
            'global minima', 'saddle point'
        ],
        'relationships': ['training', 'convergence', 'regularization'],
        'definitions': {
            'gradient descent': 'Optimization algorithm that minimizes loss by updating parameters',
            'Adam': 'Popular optimization algorithm combining momentum and adaptive learning rates',
            'learning rate': 'Step size for parameter updates during optimization',
            'regularization': 'Techniques to prevent overfitting during training'
        }
    },
    'supervised_learning': {
        'keywords': [
            'supervised learning', 'supervised', 'classification', 'regression',
            'labeled data', 'training data', 'ground truth', 'target variable',
            'predictor', 'feature', 'attribute', 'class label', 'output variable',
            'binary classification', 'multiclass', 'logistic regression',
            'decision tree', 'random forest', 'SVM', 'support vector machine',
            'gradient boosting', 'xgboost', 'lightgbm', 'cross validation',
            'train test split', 'validation set', 'model evaluation'
        ],
        'relationships': ['prediction', 'training', 'validation', 'testing'],
        'definitions': {
            'supervised learning': 'Learning paradigm where the model learns from labeled training data',
            'classification': 'Task of categorizing input data into predefined classes',
            'regression': 'Task of predicting continuous numerical values'
        }
    },
    'unsupervised_learning': {
        'keywords': [
            'unsupervised learning', 'unsupervised', 'clustering', 'dimensionality reduction',
            'k-means', 'hierarchical clustering', 'DBSCAN', 'density-based clustering',
            'PCA', 'principal component analysis', 'autoencoder', 'embedding',
            'feature extraction', 'pattern recognition', 'anomaly detection',
            'outlier detection', 'density estimation', 'gaussian mixture',
            't-SNE', 'UMAP', 'latent space', 'manifold learning'
        ],
        'relationships': ['pattern discovery', 'feature learning', 'clustering'],
        'definitions': {
            'unsupervised learning': 'Learning paradigm where the model finds patterns in unlabeled data',
            'clustering': 'Grouping similar data points together without predefined labels',
            'dimensionality reduction': 'Reducing the number of features while preserving important information'
        }
    },
    'model_optimization': {
        'keywords': [
            'optimization', 'hyperparameter tuning', 'gradient descent',
            'learning rate', 'batch size', 'epoch', 'iteration',
            'loss function', 'cost function', 'objective function',
            'regularization', 'dropout', 'batch normalization',
            'momentum', 'adam', 'rmsprop', 'sgd', 'stochastic gradient descent',
            'learning rate schedule', 'early stopping', 'model checkpoint',
            'weight decay', 'l1 regularization', 'l2 regularization'
        ],
        'relationships': ['training', 'convergence', 'performance improvement'],
        'definitions': {
            'optimization': 'Process of adjusting model parameters to minimize the loss function',
            'hyperparameter': 'Configuration variables that control the learning process',
            'regularization': 'Techniques to prevent overfitting by adding constraints to the model'
        }
    },
    'evaluation_metrics': {
        'keywords': [
            'accuracy', 'precision', 'recall', 'f1 score', 'f1-score',
            'confusion matrix', 'roc curve', 'auc', 'area under curve',
            'mean squared error', 'mse', 'rmse', 'mae', 'mean absolute error',
            'cross validation', 'k-fold', 'validation set', 'test set',
            'overfitting', 'underfitting', 'bias', 'variance', 'generalization',
            'performance metric', 'evaluation metric', 'model selection'
        ],
        'relationships': ['model assessment', 'performance measurement', 'validation'],
        'definitions': {
            'accuracy': 'Proportion of correct predictions among all predictions',
            'precision': 'Proportion of true positive predictions among all positive predictions',
            'recall': 'Proportion of true positive predictions among all actual positives'
        }
    },
    'data_preprocessing': {
        'keywords': [
            'preprocessing', 'feature engineering', 'feature selection',
            'data cleaning', 'normalization', 'standardization', 'scaling',
            'missing value', 'imputation', 'outlier detection', 'outlier removal',
            'feature transformation', 'one-hot encoding', 'label encoding',
            'dimensionality reduction', 'feature extraction', 'data augmentation',
            'data transformation', 'binning', 'sampling', 'resampling'
        ],
        'relationships': ['data preparation', 'feature creation', 'data quality'],
        'definitions': {
            'preprocessing': 'Steps to prepare raw data for machine learning algorithms',
            'feature engineering': 'Process of creating new features from existing data',
            'data cleaning': 'Process of fixing or removing incorrect, corrupted, or irrelevant data'
        }
    },
    'reinforcement_learning': {
        'keywords': [
            'reinforcement learning', 'RL', 'Q-learning', 'deep Q-network', 'DQN',
            'policy gradient', 'actor-critic', 'DDPG', 'PPO', 'A3C',
            'reward function', 'state space', 'action space', 'environment',
            'exploration', 'exploitation', 'epsilon-greedy', 'SARSA',
            'monte carlo', 'temporal difference', 'value function', 'policy',
            'markov decision process', 'MDP', 'bellman equation'
        ],
        'relationships': ['learning', 'optimization', 'decision making', 'control'],
        'definitions': {
            'reinforcement learning': 'Learning paradigm where an agent learns to make decisions by interacting with an environment',
            'Q-learning': 'Algorithm that learns to make optimal decisions by estimating action-value functions',
            'policy gradient': 'Method that directly optimizes the policy by gradient ascent on expected rewards',
            'actor-critic': 'Architecture combining value function estimation and policy optimization'
        }
    },
    'natural_language_processing': {
        'keywords': [
            'NLP', 'natural language processing', 'word embedding', 'word2vec',
            'GloVe', 'FastText', 'tokenization', 'lemmatization', 'stemming',
            'named entity recognition', 'NER', 'part of speech', 'POS tagging',
            'sentiment analysis', 'text classification', 'machine translation',
            'sequence-to-sequence', 'language model', 'BERT', 'GPT', 'T5',
            'RoBERTa', 'XLNet', 'attention mechanism', 'transformer'
        ],
        'relationships': ['text processing', 'language understanding', 'generation'],
        'definitions': {
            'NLP': 'Field focused on enabling computers to understand and process human language',
            'word embedding': 'Dense vector representations of words capturing semantic meaning',
            'tokenization': 'Process of breaking text into smaller units like words or subwords',
            'language model': 'Model that predicts the probability of sequences of words'
        }
    },
    'computer_vision': {
        'keywords': [
            'computer vision', 'CV', 'image processing', 'object detection',
            'image segmentation', 'semantic segmentation', 'instance segmentation',
            'YOLO', 'R-CNN', 'Fast R-CNN', 'Faster R-CNN', 'Mask R-CNN',
            'feature extraction', 'edge detection', 'corner detection', 'SIFT',
            'SURF', 'HOG', 'optical flow', 'pose estimation', 'face recognition',
            'image classification', 'transfer learning', 'data augmentation'
        ],
        'relationships': ['vision', 'detection', 'recognition', 'segmentation'],
        'definitions': {
            'computer vision': 'Field focused on enabling computers to understand and process visual information',
            'object detection': 'Task of locating and classifying objects in images',
            'image segmentation': 'Task of partitioning images into meaningful parts',
            'feature extraction': 'Process of identifying distinctive visual patterns'
        }
    },
    'clustering_algorithms': {
        'keywords': [
            'clustering', 'K-means', 'hierarchical clustering', 'DBSCAN',
            'mean shift', 'spectral clustering', 'agglomerative', 'divisive',
            'dendogram', 'silhouette score', 'elbow method', 'cluster centroid',
            'density-based', 'connectivity-based', 'gaussian mixture', 'GMM',
            'fuzzy c-means', 'partition-based clustering', 'cluster analysis',
            'cluster validation', 'inertia', 'cophenetic correlation'
        ],
        'relationships': ['grouping', 'segmentation', 'pattern discovery'],
        'definitions': {
            'clustering': 'Task of grouping similar data points together without supervision',
            'K-means': 'Algorithm that partitions data into K clusters based on centroids',
            'DBSCAN': 'Density-based clustering algorithm that can find arbitrary-shaped clusters',
            'hierarchical clustering': 'Method that creates a tree of clusters'
        }
    },
    'dimensionality_reduction': {
        'keywords': [
            'dimensionality reduction', 'PCA', 'principal component analysis',
            't-SNE', 'UMAP', 'LDA', 'linear discriminant analysis',
            'feature selection', 'feature extraction', 'manifold learning',
            'autoencoder', 'variational autoencoder', 'VAE', 'kernel PCA',
            'factor analysis', 'matrix factorization', 'SVD', 'NMF',
            'locally linear embedding', 'LLE', 'isomap', 'MDS'
        ],
        'relationships': ['visualization', 'compression', 'feature learning'],
        'definitions': {
            'dimensionality reduction': 'Techniques to reduce data dimensions while preserving important information',
            'PCA': 'Method that finds principal components maximizing variance',
            't-SNE': 'Non-linear technique for visualizing high-dimensional data',
            'autoencoder': 'Neural network architecture for learning compressed representations'
        }
    },
    'optimization_algorithms': {
        'keywords': [
            'optimization', 'gradient descent', 'SGD', 'Adam', 'RMSprop',
            'AdaGrad', 'AdaDelta', 'momentum', 'Nesterov', 'learning rate',
            'batch size', 'epoch', 'convergence', 'local minima', 'global minima',
            'saddle point', 'line search', 'newton method', 'quasi-newton',
            'BFGS', 'L-BFGS', 'conjugate gradient', 'particle swarm'
        ],
        'relationships': ['minimization', 'maximization', 'parameter tuning'],
        'definitions': {
            'optimization': 'Process of finding the best parameters to minimize or maximize an objective',
            'gradient descent': 'First-order optimization algorithm using gradients',
            'Adam': 'Adaptive moment estimation optimizer combining momentum and RMSprop',
            'learning rate': 'Step size parameter controlling optimization speed'
        }
    },
    'bayesian_methods': {
        'keywords': [
            'bayesian', 'bayes theorem', 'prior', 'likelihood', 'posterior',
            'bayesian inference', 'bayesian network', 'probabilistic graphical model',
            'MCMC', 'markov chain monte carlo', 'gibbs sampling', 'variational inference',
            'bayesian optimization', 'gaussian process', 'GP', 'beta distribution',
            'dirichlet process', 'hierarchical model', 'conjugate prior',
            'bayesian deep learning', 'bayesian neural network'
        ],
        'relationships': ['inference', 'probability', 'uncertainty'],
        'definitions': {
            'bayesian inference': 'Statistical inference using Bayes theorem to update probabilities',
            'bayesian network': 'Probabilistic graphical model representing conditional dependencies',
            'MCMC': 'Method for sampling from probability distributions',
            'gaussian process': 'Probabilistic approach to regression and classification'
        }
    },
    'anomaly_detection': {
        'keywords': [
            'anomaly detection', 'outlier detection', 'novelty detection',
            'one-class SVM', 'isolation forest', 'local outlier factor', 'LOF',
            'autoencoder', 'reconstruction error', 'statistical testing',
            'density estimation', 'distance-based', 'clustering-based',
            'ensemble methods', 'robust statistics', 'mahalanobis distance',
            'time series anomaly', 'fraud detection', 'system monitoring'
        ],
        'relationships': ['detection', 'monitoring', 'security'],
        'definitions': {
            'anomaly detection': 'Task of identifying unusual patterns that deviate from expected behavior',
            'isolation forest': 'Algorithm that isolates anomalies using random splits',
            'local outlier factor': 'Method that finds anomalies by comparing local densities',
            'one-class SVM': 'Support vector method for novelty detection'
        }
    },
    'time_series_analysis': {
        'keywords': [
            'time series', 'forecasting', 'ARIMA', 'SARIMA', 'exponential smoothing',
            'prophet', 'RNN', 'LSTM', 'GRU', 'sequence modeling', 'seasonality',
            'trend', 'stationarity', 'autocorrelation', 'moving average',
            'decomposition', 'fourier transform', 'wavelet transform',
            'change point detection', 'dynamic time warping', 'DTW'
        ],
        'relationships': ['prediction', 'forecasting', 'sequence analysis'],
        'definitions': {
            'time series': 'Data points indexed in time order',
            'ARIMA': 'AutoRegressive Integrated Moving Average model for forecasting',
            'prophet': 'Forecasting tool that handles daily data with strong seasonality',
            'dynamic time warping': 'Algorithm for measuring similarity between temporal sequences'
        }
    },
    'meta_learning': {
        'keywords': [
            'meta learning', 'few-shot learning', 'one-shot learning',
            'zero-shot learning', 'transfer learning', 'multi-task learning',
            'MAML', 'prototypical networks', 'matching networks',
            'learning to learn', 'meta-gradients', 'hyperparameter optimization',
            'neural architecture search', 'NAS', 'autoML', 'model selection',
            'cross-validation', 'ensemble selection', 'knowledge distillation'
        ],
        'relationships': ['adaptation', 'generalization', 'transfer'],
        'definitions': {
            'meta learning': 'Learning to learn, or improving learning ability over multiple tasks',
            'few-shot learning': 'Learning from very few examples per class',
            'transfer learning': 'Applying knowledge from one task to another',
            'autoML': 'Automated machine learning pipeline optimization'
        }
    },
    'support_vector_machines': {
        'keywords': [
            'SVM', 'support vector machine', 'kernel trick', 'kernel function',
            'linear kernel', 'RBF kernel', 'polynomial kernel', 'margin',
            'support vectors', 'hyperplane', 'soft margin', 'hard margin',
            'C parameter', 'gamma parameter', 'dual formulation', 'primal formulation',
            'SMO algorithm', 'quadratic programming', 'slack variables'
        ],
        'relationships': ['classification', 'regression', 'feature transformation'],
        'definitions': {
            'SVM': 'Machine learning algorithm that finds the optimal hyperplane for classification or regression',
            'kernel trick': 'Method to perform calculations in higher dimensional space without explicit transformation',
            'support vectors': 'Data points that lie closest to the decision boundary and define the margin',
            'margin': 'Distance between the decision boundary and the nearest data points'
        },
        'practical_cases': {
            'text_classification': {
                'scenario': 'Text classification system for spam detection',
                'implementation': 'Using TF-IDF features with RBF kernel SVM',
                'evaluation': {'accuracy': '98%', 'precision': '97%', 'recall': '96%'},
                'challenges': 'Handling imbalanced classes and feature selection'
            },
            'bioinformatics': {
                'scenario': 'Protein function prediction from sequence data',
                'implementation': 'Custom string kernel SVM for sequence analysis',
                'evaluation': {'accuracy': '92%', 'specificity': '94%'},
                'challenges': 'Designing appropriate kernel for biological sequences'
            }
        },
        'numerical_examples': {
            'kernel_computation': {
                'scenario': 'Given two vectors x1=[1,2] and x2=[3,4], compute RBF kernel value',
                'calculation': 'K(x1,x2) = exp(-gamma * ||x1-x2||^2)',
                'result': 'With gamma=0.1, K(x1,x2) = exp(-0.1 * 8) ≈ 0.45'
            }
        }
    },
    'recurrent_neural_networks': {
        'keywords': [
            'RNN', 'recurrent neural network', 'LSTM', 'GRU', 'bidirectional RNN',
            'sequence modeling', 'backpropagation through time', 'BPTT',
            'vanishing gradient', 'exploding gradient', 'memory cell', 'gates',
            'hidden state', 'cell state', 'sequence-to-sequence', 'encoder-decoder',
            'attention mechanism', 'teacher forcing', 'truncated BPTT'
        ],
        'relationships': ['sequence learning', 'time series', 'language modeling'],
        'definitions': {
            'RNN': 'Neural network architecture designed to process sequential data',
            'LSTM': 'Long Short-Term Memory network that addresses vanishing gradient problem',
            'GRU': 'Gated Recurrent Unit, a simplified version of LSTM',
            'BPTT': 'Algorithm for training RNNs by unrolling the network through time'
        },
        'practical_cases': {
            'language_model': {
                'scenario': 'Building a language model for code completion',
                'implementation': 'Stacked LSTM with attention mechanism',
                'evaluation': {'perplexity': '32.5', 'accuracy': '85%'},
                'challenges': 'Handling long-range dependencies and vocabulary size'
            },
            'stock_prediction': {
                'scenario': 'Time series forecasting for stock prices',
                'implementation': 'Bidirectional GRU with technical indicators',
                'evaluation': {'RMSE': '0.023', 'MAE': '0.018'},
                'challenges': 'Dealing with market volatility and feature engineering'
            }
        },
        'numerical_examples': {
            'gradient_calculation': {
                'scenario': 'Calculate gradients for a simple RNN with tanh activation',
                'calculation': 'dh_t/dW = tanh\'(Wx_t + Uh_{t-1}) * x_t',
                'result': 'Numerical example with specific weight matrix and input sequence'
            }
        }
    },
    'hidden_markov_models': {
        'keywords': [
            'HMM', 'hidden markov model', 'markov chain', 'emission probability',
            'transition probability', 'hidden states', 'observable states',
            'forward algorithm', 'backward algorithm', 'viterbi algorithm',
            'baum-welch algorithm', 'EM for HMM', 'state sequence',
            'maximum likelihood', 'forward-backward algorithm'
        ],
        'relationships': ['sequence modeling', 'probabilistic modeling', 'state estimation'],
        'definitions': {
            'HMM': 'Statistical model where system is assumed to be a Markov process with hidden states',
            'viterbi algorithm': 'Dynamic programming algorithm to find most likely sequence of hidden states',
            'baum-welch': 'EM algorithm variant for learning HMM parameters',
            'forward algorithm': 'Method to compute probability of observation sequence'
        },
        'practical_cases': {
            'speech_recognition': {
                'scenario': 'Phoneme recognition system using HMM',
                'implementation': 'Multi-state HMM with MFCC features',
                'evaluation': {'accuracy': '89%', 'error_rate': '11%'},
                'challenges': 'Handling variable length utterances and noise'
            }
        },
        'numerical_examples': {
            'probability_computation': {
                'scenario': 'Given a 2-state HMM, compute probability of observation sequence',
                'calculation': 'Using forward algorithm with transition and emission matrices',
                'result': 'Step-by-step probability calculation for sequence [0,1,1]'
            }
        }
    },
    'expectation_maximization': {
        'keywords': [
            'EM algorithm', 'expectation maximization', 'latent variables',
            'maximum likelihood', 'E-step', 'M-step', 'convergence',
            'gaussian mixture model', 'GMM', 'parameter estimation',
            'incomplete data', 'likelihood function', 'jensen inequality',
            'local optimum', 'initialization strategies'
        ],
        'relationships': ['clustering', 'density estimation', 'parameter learning'],
        'definitions': {
            'EM algorithm': 'Iterative method to find maximum likelihood parameters in latent variable models',
            'E-step': 'Expectation step that computes expected values of latent variables',
            'M-step': 'Maximization step that updates parameters to maximize likelihood',
            'GMM': 'Probabilistic model that assumes data comes from mixture of Gaussians'
        },
        'practical_cases': {
            'image_segmentation': {
                'scenario': 'Image segmentation using Gaussian Mixture Models',
                'implementation': 'EM for learning GMM parameters from pixel intensities',
                'evaluation': {'accuracy': '91%', 'IoU': '0.85'},
                'challenges': 'Choosing number of components and initialization'
            }
        },
        'numerical_examples': {
            'gmm_update': {
                'scenario': 'Update GMM parameters for 2-component mixture',
                'calculation': 'E-step: compute responsibilities, M-step: update means and covariances',
                'result': 'Numerical example with 2D data points'
            }
        }
    },
    'map_and_maximum_likelihood': {
        'keywords': [
            'MAP', 'maximum a posteriori', 'MLE', 'maximum likelihood estimation',
            'prior probability', 'posterior probability', 'likelihood function',
            'bayesian inference', 'point estimation', 'conjugate prior',
            'regularization', 'parameter estimation', 'optimization',
            'log likelihood', 'probability theory'
        ],
        'relationships': ['estimation', 'inference', 'optimization'],
        'definitions': {
            'MAP': 'Maximum a posteriori estimation incorporating prior knowledge',
            'MLE': 'Maximum likelihood estimation based only on observed data',
            'prior probability': 'Initial belief about parameters before observing data',
            'posterior probability': 'Updated belief about parameters after observing data'
        },
        'practical_cases': {
            'regression': {
                'scenario': 'Linear regression with MAP estimation',
                'implementation': 'Using Gaussian prior for regularization',
                'evaluation': {'MSE': '0.15', 'R2': '0.89'},
                'challenges': 'Choosing appropriate prior distribution'
            }
        },
        'numerical_examples': {
            'coin_flip': {
                'scenario': 'Estimate probability of heads from coin flips using MAP',
                'calculation': 'MAP = (heads + alpha)/(total + alpha + beta)',
                'result': 'With 7 heads in 10 flips and Beta(2,2) prior'
            }
        }
    },
    'advanced_gan_architectures': {
        'keywords': [
            'DCGAN', 'WGAN', 'conditional GAN', 'CycleGAN', 'StyleGAN',
            'progressive GAN', 'BigGAN', 'SAGAN', 'InfoGAN', 'BEGAN',
            'adversarial training', 'generator loss', 'discriminator loss',
            'mode collapse', 'gradient penalty', 'spectral normalization',
            'style mixing', 'adaptive instance normalization', 'AdaIN'
        ],
        'relationships': ['generation', 'style transfer', 'domain adaptation'],
        'definitions': {
            'DCGAN': 'Deep Convolutional GAN using convolutional architectures',
            'WGAN': 'Wasserstein GAN using earth mover distance',
            'StyleGAN': 'GAN architecture for controllable high-quality image generation',
            'CycleGAN': 'Unpaired image-to-image translation using cycle consistency'
        },
        'practical_cases': {
            'face_generation': {
                'scenario': 'Photorealistic face generation system',
                'implementation': 'StyleGAN2 with custom dataset',
                'evaluation': {'FID': '4.40', 'IS': '8.52'},
                'challenges': 'Training stability and diversity'
            },
            'style_transfer': {
                'scenario': 'Artistic style transfer application',
                'implementation': 'CycleGAN with attention mechanism',
                'evaluation': {'user_satisfaction': '92%', 'style_accuracy': '88%'},
                'challenges': 'Preserving content while transferring style'
            }
        },
        'numerical_examples': {
            'loss_computation': {
                'scenario': 'Compute WGAN loss with gradient penalty',
                'calculation': 'L = E[D(fake)] - E[D(real)] + λE[(||∇D(x)||₂ - 1)²]',
                'result': 'Numerical example with batch of images'
            }
        }
    },
    'supervised_learning_detailed': {
        'keywords': [
            'supervised learning', 'classification', 'regression', 'labeled data',
            'training data', 'test data', 'validation set', 'cross validation',
            'overfitting', 'underfitting', 'bias-variance tradeoff', 'regularization',
            'feature selection', 'feature engineering', 'model evaluation',
            'confusion matrix', 'precision', 'recall', 'f1-score', 'accuracy',
            'ROC curve', 'AUC', 'loss function', 'gradient descent'
        ],
        'relationships': ['training', 'prediction', 'evaluation', 'optimization'],
        'definitions': {
            'supervised learning': 'Learning paradigm where the model learns from labeled examples to make predictions',
            'classification': 'Task of categorizing input data into predefined classes',
            'regression': 'Task of predicting continuous numerical values',
            'cross validation': 'Technique to assess model performance on different data splits'
        },
        'practical_cases': {
            'spam_detection': {
                'scenario': 'Email spam classification system',
                'implementation': 'Using text features with logistic regression',
                'evaluation': {'accuracy': '97%', 'precision': '95%', 'recall': '96%'},
                'challenges': 'Handling imbalanced classes and feature extraction'
            },
            'house_price': {
                'scenario': 'Real estate price prediction',
                'implementation': 'Multiple regression with regularization',
                'evaluation': {'RMSE': '45000', 'R2': '0.89'},
                'challenges': 'Feature selection and handling outliers'
            }
        },
        'numerical_examples': {
            'cross_validation': {
                'scenario': 'Implement 5-fold cross validation on dataset with 1000 samples',
                'calculation': 'Split size = 1000/5 = 200 samples per fold',
                'result': 'Train on 800 samples, validate on 200, rotate 5 times'
            }
        }
    },
    'unsupervised_learning_detailed': {
        'keywords': [
            'unsupervised learning', 'clustering', 'dimensionality reduction',
            'density estimation', 'anomaly detection', 'pattern discovery',
            'k-means', 'hierarchical clustering', 'DBSCAN', 'mean shift',
            'PCA', 't-SNE', 'UMAP', 'autoencoders', 'self-organizing maps',
            'gaussian mixture models', 'latent variables', 'manifold learning',
            'feature learning', 'representation learning'
        ],
        'relationships': ['clustering', 'visualization', 'pattern discovery', 'feature learning'],
        'definitions': {
            'unsupervised learning': 'Learning paradigm that finds patterns in unlabeled data',
            'clustering': 'Grouping similar data points without predefined labels',
            'dimensionality reduction': 'Reducing data dimensions while preserving important information',
            'anomaly detection': 'Finding unusual patterns that deviate from expected behavior'
        },
        'practical_cases': {
            'customer_segmentation': {
                'scenario': 'Customer segmentation for targeted marketing',
                'implementation': 'K-means clustering on customer behavior data',
                'evaluation': {'silhouette_score': '0.75', 'calinski_score': '156.8'},
                'challenges': 'Determining optimal number of clusters'
            },
            'image_compression': {
                'scenario': 'Image compression using PCA',
                'implementation': 'PCA with variance retention threshold',
                'evaluation': {'variance_retained': '95%', 'compression_ratio': '4:1'},
                'challenges': 'Balancing compression ratio and image quality'
            }
        },
        'numerical_examples': {
            'kmeans_calculation': {
                'scenario': 'Calculate cluster centroids for 2D data points',
                'calculation': 'Mean of points: centroid = (Σx_i/n, Σy_i/n)',
                'result': 'Example with specific coordinate calculations'
            }
        }
    },
    'probability_and_statistics': {
        'keywords': [
            'probability theory', 'statistics', 'random variables', 'distributions',
            'mean', 'variance', 'standard deviation', 'correlation', 'covariance',
            'hypothesis testing', 'p-value', 'confidence interval', 'likelihood',
            'bayesian statistics', 'frequentist statistics', 'sampling',
            'central limit theorem', 'law of large numbers', 'probability density',
            'joint probability', 'conditional probability', 'bayes theorem'
        ],
        'relationships': ['inference', 'estimation', 'testing', 'modeling'],
        'definitions': {
            'probability': 'Mathematical framework for modeling uncertainty and randomness',
            'statistics': 'Methods for collecting, analyzing, and drawing conclusions from data',
            'hypothesis testing': 'Framework for making decisions based on data',
            'bayesian inference': 'Updating beliefs based on new evidence using Bayes theorem'
        },
        'practical_cases': {
            'ab_testing': {
                'scenario': 'A/B testing for website conversion rate',
                'implementation': 'Chi-square test with significance level 0.05',
                'evaluation': {'p_value': '0.023', 'effect_size': '0.15'},
                'challenges': 'Determining sample size and handling multiple testing'
            },
            'medical_diagnosis': {
                'scenario': 'Disease diagnosis using Bayesian inference',
                'implementation': 'Naive Bayes with multiple symptoms',
                'evaluation': {'sensitivity': '92%', 'specificity': '95%'},
                'challenges': 'Handling conditional independence assumptions'
            }
        },
        'numerical_examples': {
            'bayes_theorem': {
                'scenario': 'Calculate posterior probability for medical test',
                'calculation': 'P(D|+) = P(+|D)P(D)/P(+)',
                'result': 'With P(D)=0.01, P(+|D)=0.95, P(+|not D)=0.10'
            },
            'confidence_interval': {
                'scenario': 'Calculate 95% confidence interval for mean',
                'calculation': 'CI = x̄ ± (1.96 * σ/√n)',
                'result': 'Example with sample mean, std dev, and sample size'
            }
        }
    }
}

# Bloom's Taxonomy levels and question templates
BLOOMS_TAXONOMY = {
    'remember': {
        'level': 1,
        'keywords': ['define', 'identify', 'list', 'name', 'recall', 'state', 'recognize', 'select'],
        'question_stems': [
            "What is the definition of {}?",
            "List the key components of {}",
            "Identify the main characteristics of {}",
            "State the purpose of {}"
        ],
        'mcq_templates': [
            "Which of the following best defines {}?",
            "What is the primary purpose of {}?",
            "Which statement correctly identifies {}?"
        ],
        'descriptive_templates': [
            "Define {} and list its key characteristics.",
            "State the main components of {} and their functions."
        ]
    },
    'understand': {
        'level': 2,
        'keywords': ['explain', 'compare', 'contrast', 'discuss', 'analyze', 'interpret', 'summarize', 'describe'],
        'question_stems': [
            "Explain the differences between {} and {}",
            "Compare and contrast {} and {}",
            "Discuss the similarities and differences between {} and {}",
            "Analyze the role of {} in {}"
        ],
        'mcq_templates': [
            "Which of the following best describes {}?",
            "What is the main purpose of {}?",
            "Which statement correctly identifies {}?"
        ],
        'descriptive_templates': [
            "Describe {} and explain its significance.",
            "Discuss the key features of {} and their implications.",
            "Explain how {} works and its importance in {}."
        ]
    },
    'apply': {
        'level': 3,
        'keywords': ['implement', 'use', 'apply', 'adapt', 'modify', 'combine', 'integrate', 'apply'],
        'question_stems': [
            "How can {} be implemented in a practical context?",
            "What are the steps to apply {} effectively?",
            "How can {} be adapted to different situations?",
            "What are the implications of using {} in {}?"
        ],
        'mcq_templates': [
            "Which of the following best describes how to apply {}?",
            "What is the most effective way to use {}?",
            "How can {} be modified to fit different contexts?",
            "What are the consequences of using {} in {}?"
        ],
        'descriptive_templates': [
            "Discuss the practical application of {} in a real-world scenario.",
            "Explain how {} can be applied in different situations.",
            "Describe the process of implementing {} and its impact on {}."
        ]
    },
    'analyze': {
        'level': 4,
        'keywords': ['evaluate', 'critique', 'compare', 'contrast', 'break down', 'analyze', 'interpret', 'evaluate'],
        'question_stems': [
            "Evaluate the effectiveness of {} in a real-world scenario.",
            "Critique the strengths and weaknesses of {}.",
            "Compare and contrast {} with alternatives.",
            "Break down the components of {} and explain their roles."
        ],
        'mcq_templates': [
            "Which of the following best describes the strengths and weaknesses of {}?",
            "What are the key factors to consider when evaluating {}?",
            "How does {} compare to alternatives in terms of effectiveness?",
            "What are the implications of using {} in a real-world scenario?"
        ],
        'descriptive_templates': [
            "Analyze the strengths and weaknesses of {} in a detailed manner.",
            "Discuss the factors that contribute to the effectiveness of {}.",
            "Break down the components of {} and explain their roles in a real-world context."
        ]
    },
    'evaluate': {
        'level': 5,
        'keywords': ['judge', 'assess', 'compare', 'evaluate', 'justify', 'evaluate', 'compare', 'assess'],
        'question_stems': [
            "How can {} be justified in a real-world scenario?",
            "Evaluate the benefits and drawbacks of {}.",
            "Compare and contrast {} with alternatives.",
            "Assess the effectiveness of {} in a real-world context."
        ],
        'mcq_templates': [
            "Which of the following best describes the benefits and drawbacks of {}?",
            "What are the key factors to consider when evaluating {}?",
            "How does {} compare to alternatives in terms of effectiveness?",
            "What are the implications of using {} in a real-world scenario?"
        ],
        'descriptive_templates': [
            "Justify the use of {} in a real-world scenario.",
            "Discuss the benefits and drawbacks of {} in detail.",
            "Compare and contrast {} with alternatives and explain why {} is the best choice."
        ]
    },
    'create': {
        'level': 6,
        'keywords': ['design', 'develop', 'create', 'innovate', 'invent', 'design', 'develop', 'create'],
        'question_stems': [
            "How can {} be designed to solve a specific problem?",
            "What are the steps to develop {} effectively?",
            "How can {} be innovated to create a new solution?",
            "What are the implications of designing {} for a real-world application?"
        ],
        'mcq_templates': [
            "Which of the following best describes how to design {}?",
            "What is the most effective way to develop {}?",
            "How can {} be innovated to create a new solution?",
            "What are the consequences of designing {} for a real-world application?"
        ],
        'descriptive_templates': [
            "Discuss the design process of {} and its implications.",
            "Explain how {} can be developed effectively.",
            "Describe the steps to create {} and its impact on a real-world application."
        ]
    }
}

# Global metrics tracking
METRICS = {
    'question_generation_times': [],
    'classification_results': {'correct': 0, 'total': 0},
    'teacher_validations': {'approved': 0, 'total': 0},
    'start_time': datetime.now(),
    'downtime_periods': []
}

def calculate_metrics():
    """Calculate current performance metrics"""
    try:
        # Calculate average question generation time
        avg_gen_time = sum(METRICS['question_generation_times']) / len(METRICS['question_generation_times']) if METRICS['question_generation_times'] else 0
        
        # Calculate classification accuracy
        classification_accuracy = (METRICS['classification_results']['correct'] / METRICS['classification_results']['total'] * 100) if METRICS['classification_results']['total'] > 0 else 0
        
        # Calculate teacher validation accuracy
        teacher_accuracy = (METRICS['teacher_validations']['approved'] / METRICS['teacher_validations']['total'] * 100) if METRICS['teacher_validations']['total'] > 0 else 0
        
        # Calculate uptime
        total_duration = datetime.now() - METRICS['start_time']
        downtime = sum((end - start).total_seconds() for start, end in METRICS['downtime_periods'])
        uptime_percentage = ((total_duration.total_seconds() - downtime) / total_duration.total_seconds() * 100) if total_duration.total_seconds() > 0 else 100
        
        return {
            'avg_question_generation_time': f"{avg_gen_time:.2f}s",
            'classification_accuracy': f"{classification_accuracy:.1f}%",
            'teacher_validation_accuracy': f"{teacher_accuracy:.1f}%",
            'backend_uptime': f"{uptime_percentage:.2f}%",
            'total_questions_generated': METRICS['classification_results']['total'],
            'total_teacher_validations': METRICS['teacher_validations']['total'],
            'system_start_time': METRICS['start_time'].strftime('%Y-%m-%d %H:%M:%S'),
            'uptime_duration': str(total_duration).split('.')[0]
        }
    except Exception as e:
        logger.error(f"Error calculating metrics: {str(e)}")
        return None

@app.route('/metrics', methods=['GET'])
def get_metrics():
    """Get current system performance metrics"""
    try:
        metrics = calculate_metrics()
        if not metrics:
            return jsonify({
                'success': False,
                'error': 'Error calculating metrics'
            }), 500
            
        return jsonify({
            'success': True,
            'metrics': metrics
        })
    except Exception as e:
        logger.error(f"Error in metrics endpoint: {str(e)}")
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

def initialize_model():
    """Initialize the base model and tokenizer"""
    global base_model, tokenizer
    try:
        logger.info("Loading base model and tokenizer...")
        model_name = "gpt2-medium"  # or another base model of your choice
        base_model = AutoModelForCausalLM.from_pretrained(model_name)
        tokenizer = AutoTokenizer.from_pretrained(model_name)
        
        # Move model to GPU if available
        device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
        base_model = base_model.to(device)
        
        logger.info(f"Base model loaded successfully on {device}")
        return True
    except Exception as e:
        logger.error(f"Model loading failed: {str(e)}")
        return False

def create_peft_model(content):
    """Create a PEFT model adapted to the current content"""
    global base_model, peft_model, current_context_hash, tokenizer
    try:
        new_context_hash = hash(content) % 10000
        logger.info(f"New context hash: {new_context_hash}")
        
        # Extract concepts to adapt the model
        detected_concepts = extract_ml_concepts(content)
        logger.info(f"Detected concepts for PEFT adaptation: {list(detected_concepts.keys())}")
        
        # Configure LoRA with compatible target modules
        lora_config = LoraConfig(
            task_type=TaskType.CAUSAL_LM,
            r=16,  # Reduced rank for stability
            lora_alpha=32,
            lora_dropout=0.1,
            bias="none",
            target_modules=["c_attn", "c_proj"],  # Only target attention modules
            inference_mode=False
        )
        
        # Create PEFT model with error handling
        try:
            peft_model = get_peft_model(base_model, lora_config)
            logger.info("PEFT model created successfully")
        except Exception as e:
            logger.error(f"Error creating PEFT model: {str(e)}")
            # Fallback to base model if PEFT creation fails
            peft_model = base_model
            logger.info("Falling back to base model")
        
        # Create adaptation text
        adaptation_text = []
        
        # Add topic markers
        adaptation_text.extend([
            "CONTENT CONTEXT:",
            "===============",
            content[:500],  # Include start of content
            "===============",
            ""
        ])
        
        # Add concept information
        for topic, data in detected_concepts.items():
            adaptation_text.extend([
                f"TOPIC: {topic.upper()}",
                "----------------",
                f"Main Concepts: {', '.join(data['keywords'][:5])}",
                "",
                "Key Definitions:",
                *[f"- {k}: {v}" for k, v in list(data['definitions'].items())[:3]],
                "",
                "----------------",
                ""
            ])
        
        # Encode and process adaptation text
        logger.info("Encoding adaptation text...")
        inputs = tokenizer("\n".join(adaptation_text), return_tensors="pt", truncation=True, max_length=1024)
        inputs = {k: v.to(peft_model.device) for k, v in inputs.items()}
        
        # Adapt model
        logger.info("Adapting model...")
        with torch.no_grad():
            outputs = peft_model(**inputs)
        
        current_context_hash = new_context_hash
        logger.info("Model adaptation completed successfully")
        return True
        
    except Exception as e:
        logger.error(f"Error in model creation: {str(e)}")
        return False

def extract_text_from_docx(filepath):
    """Extract text from a .docx file"""
    try:
        doc = Document(filepath)
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text.append(cell.text)
        return '\n'.join(text)
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {str(e)}")
        raise

def extract_text_from_pptx(filepath):
    """Extract text from a .pptx file"""
    try:
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
            # Extract text from notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                for shape in notes_slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {str(e)}")
        raise

def extract_text_from_file(filepath):
    """Extract text from various file types"""
    try:
        logger.info(f"Starting text extraction from: {filepath}")
        ext = os.path.splitext(filepath)[1].lower()
        
        # Clear previous content files
        summary_dir = Path(app.config['SUMMARIES_DIR'])
        if summary_dir.exists():
            for f in summary_dir.glob('*.txt'):
                try:
                    f.unlink()
                    logger.info(f"Removed old file: {f}")
                except Exception as e:
                    logger.error(f"Error removing old file {f}: {str(e)}")
        
        if ext == '.pdf':
            logger.info("Extracting text from PDF")
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = []
                for i, page in enumerate(reader.pages):
                    logger.info(f"Processing PDF page {i+1}/{len(reader.pages)}")
                    page_text = page.extract_text()
                    if page_text:
                        text.append(page_text)
                return '\n'.join(text)
                
        elif ext == '.txt':
            logger.info("Extracting text from TXT")
            with open(filepath, 'r', encoding='utf-8') as f:
                return f.read()
                
        elif ext == '.docx':
            logger.info("Extracting text from DOCX")
            doc = Document(filepath)
            text = []
            logger.info(f"Processing {len(doc.paragraphs)} paragraphs")
            for para in doc.paragraphs:
                if para.text:
                    text.append(para.text)
            logger.info(f"Processing {len(doc.tables)} tables")
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            text.append(cell.text)
            return '\n'.join(text)
            
        elif ext == '.pptx':
            logger.info("Extracting text from PPTX")
            prs = Presentation(filepath)
            text = []
            logger.info(f"Processing {len(prs.slides)} slides")
            for i, slide in enumerate(prs.slides):
                logger.info(f"Processing slide {i+1}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text.append(shape.text)
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    for shape in notes_slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text.append(shape.text)
            return '\n'.join(text)
            
        else:
            raise ValueError(f"Unsupported file type: {ext}")
            
    except Exception as e:
        logger.error(f"Error extracting text from {filepath}: {str(e)}")
        logger.error(traceback.format_exc())
        raise

def save_summary(filename, text):
    """Save content and summary to files"""
    try:
        # Save full content
        content_path = os.path.join(app.config['SUMMARIES_DIR'], f"{filename}.txt")
        with open(content_path, 'w', encoding='utf-8') as f:
            f.write(text)
        
        # Generate and save summary
        summary = generate_summary(text)
        summary_path = os.path.join(app.config['SUMMARIES_DIR'], f"{filename}_summary.txt")
        with open(summary_path, 'w', encoding='utf-8') as f:
            f.write(summary)
            
        return content_path, summary_path
    except Exception as e:
        logger.error(f"Error saving summary: {str(e)}")
        return None, None

def generate_key_terms(text, n=10):
    # Simple key term extraction using TF-IDF
    vectorizer = TfidfVectorizer(
        stop_words='english', 
        max_features=n,
        ngram_range=(1, 2)  # Allow for bigrams
    )
    try:
        vectorizer.fit([text])
        return vectorizer.get_feature_names_out().tolist()
    except:
        # Fallback to simple word frequency for very short texts
        words = text.lower().split()
        word_freq = defaultdict(int)
        for word in words:
            if len(word) > 3:  # Only consider words longer than 3 characters
                word_freq[word] += 1
        return sorted(word_freq, key=word_freq.get, reverse=True)[:n]

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint for Docker"""
    return jsonify({
        'status': 'healthy',
        'service': 'ml-service',
        'timestamp': datetime.now().isoformat()
    }), 200

@app.route('/')
def index():
    return jsonify({
        'status': 'running',
        'message': 'ML Question Generation API is running',
        'endpoints': {
            '/process-content': 'POST - Process uploaded content',
            '/load-model': 'POST - Load model with content',
            '/generate-questions': 'POST - Generate questions',
            '/generate-paper': 'POST - Generate paper',
            '/papers/<paper_id>': 'GET - Get paper, DELETE - Delete paper'
        }
    })

@app.route('/process-content', methods=['POST'])
def process_content():
    """Process uploaded files and generate summaries"""
    try:
        logger.info("Starting file processing...")
        
        if 'files[]' not in request.files:
            logger.error("No files found in request")
            return jsonify({'success': False, 'error': 'No files provided'}), 400

        # Get question counts from request and enforce max limit of 4
        question_counts = {
            'mcq': min(4, int(request.form.get('mcqCount', 4))),
            'descriptive': min(4, int(request.form.get('descriptiveCount', 4))),
            'true_false': min(4, int(request.form.get('trueFalseCount', 4))),
            'fill_in_blanks': min(4, int(request.form.get('fillBlanksCount', 4))),
            'case_study': min(4, int(request.form.get('caseStudyCount', 1))),
            'numerical': min(4, int(request.form.get('numericalCount', 1)))
        }
        
        logger.info(f"Requested question counts (max 4 per type): {question_counts}")

        files = request.files.getlist('files[]')
        logger.info(f"Received {len(files)} files")

        processed_files = []
        
        # Ensure directories exist
        os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
        os.makedirs(app.config['SUMMARIES_DIR'], exist_ok=True)
        
        # Clear old content and models
        global peft_model, current_context_hash
        peft_model = None
        current_context_hash = None
        
        summary_dir = Path(app.config['SUMMARIES_DIR'])
        if summary_dir.exists():
            for f in summary_dir.glob('*.txt'):
                try:
                    f.unlink()
                    logger.info(f"Removed old file: {f}")
                except Exception as e:
                    logger.error(f"Error removing old file {f}: {str(e)}")
        
        for file in files:
            try:
                if not file or not file.filename:
                    logger.warning("Empty file object received")
                    continue

                logger.info(f"Processing file: {file.filename}")
                
                if not allowed_file(file.filename):
                    logger.warning(f"File type not allowed: {file.filename}")
                    continue

                filename = secure_filename(file.filename)
                upload_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                
                logger.info(f"Saving file to: {upload_path}")
                file.save(upload_path)
                
                text = extract_text_from_file(upload_path)
                
                if not text.strip():
                    logger.warning(f"No text content found in: {filename}")
                    os.remove(upload_path)
                    continue
                
                logger.info(f"Successfully extracted text from: {filename}")
                
                # Generate summary
                summary = generate_summary(text)
                logger.info("Summary generated")
                
                # Save content and summary
                base_name = os.path.splitext(filename)[0]
                content_path = os.path.join(app.config['SUMMARIES_DIR'], f"{base_name}_content.txt")
                summary_path = os.path.join(app.config['SUMMARIES_DIR'], f"{base_name}_summary.txt")
                
                with open(content_path, 'w', encoding='utf-8') as f:
                    f.write(text)
                with open(summary_path, 'w', encoding='utf-8') as f:
                    f.write(summary)
                
                logger.info(f"Saved content to {content_path} and summary to {summary_path}")
                
                # Generate questions with specified counts
                questions = generate_questions_from_summary(summary, question_counts)
                
                processed_files.append({
                    'filename': filename,
                    'text': text,
                    'summary': summary,
                    'questions': questions if questions else {},
                    'content_path': content_path,
                    'summary_path': summary_path
                })
                
                os.remove(upload_path)
                logger.info(f"Successfully processed: {filename}")
                
            except Exception as e:
                logger.error(f"Error processing file {file.filename}: {str(e)}")
                if 'upload_path' in locals() and os.path.exists(upload_path):
                    os.remove(upload_path)
                continue

        if not processed_files:
            return jsonify({
                'success': False, 
                'error': 'No files were successfully processed'
            }), 400

        return jsonify({
            'success': True,
            'processed_files': processed_files
        })
        
    except Exception as e:
        logger.error(f"Process content error: {str(e)}")
        return jsonify({'success': False, 'error': str(e)}), 500

def load_context():
    """Load content from all files in the summaries directory"""
    try:
        context = ""
        summary_dir = Path(app.config['SUMMARIES_DIR'])
        if not summary_dir.exists():
            logger.warning(f"Summaries directory not found: {summary_dir}")
            return ""

        # Get all content files sorted by modification time
        files = list(summary_dir.glob('*.txt'))
        if not files:
            logger.warning("No files found in summaries directory")
            return ""
            
        # Get all non-summary files
        content_files = [f for f in files if not f.name.endswith('_summary.txt')]
        if not content_files:
            logger.warning("No content files found")
            return ""
            
        # Get the most recent file
        latest_file = max(content_files, key=lambda x: x.stat().st_mtime)
        
        try:
            with open(latest_file, 'r', encoding='utf-8') as f:
                context = f.read()
                logger.info(f"Loaded content from {latest_file}, length: {len(context)}")
        except Exception as e:
            logger.error(f"Error reading file {latest_file}: {str(e)}")
            return ""

        if not context:
            logger.warning("No content found in file")
        else:
            logger.info(f"Total content length: {len(context)}")
            
        return context
    except Exception as e:
        logger.error(f"Error loading context: {str(e)}")
        return ""

@app.route('/load-model', methods=['POST'])
def load_model_endpoint():
    try:
        # Check if we have content files
        summary_dir = Path(app.config['SUMMARIES_DIR'])
        if not summary_dir.exists() or not list(summary_dir.glob('*_content.txt')):
            logger.error("No content files found")
            return jsonify({
                'success': False,
                'error': 'No content available. Please upload content first.'
            }), 400

        # Get the most recent content file
        content_files = list(summary_dir.glob('*_content.txt'))
        latest_content = max(content_files, key=lambda x: x.stat().st_mtime)
        
        with open(latest_content, 'r', encoding='utf-8') as f:
            content = f.read()
            
        if not content.strip():
            logger.error("Content file is empty")
            return jsonify({
                'success': False,
                'error': 'Content file is empty'
            }), 400

        # Create PEFT model for the content
        if not create_peft_model(content):
            logger.error("Failed to create PEFT model")
            return jsonify({
                'success': False,
                'error': 'Failed to create PEFT model'
            }), 500

        # Calculate basic statistics
        word_count = len(content.split())
        sentences = [s.strip() for s in content.split('.') if s.strip()]
        
        return jsonify({
            'success': True,
            'stats': {
                'word_count': word_count,
                'sentence_count': len(sentences),
                'file_count': len(list(summary_dir.glob('*.txt')))
            },
            'sample_content': sentences[0][:100] + '...' if sentences else ''
        })
    except Exception as e:
        logger.error(f"Load model error: {str(e)}")
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

def extract_ml_concepts(text):
    """Extract ML-related concepts from text with improved pattern matching"""
    try:
        found_concepts = {}
        
        if not text or len(text.strip()) < 100:
            logger.warning("Text is too short for meaningful concept extraction")
            return {
                'general_ml': {
                    'keywords': ['machine learning'],
                    'definitions': {
                        'machine learning': 'A field of study that gives computers the ability to learn without being explicitly programmed'
                    },
                    'relationships': ['learning', 'prediction'],
                    'numerical_examples': {},
                    'practical_cases': {}
                }
            }
        
        # Preprocess text for better matching
        text_lower = text.lower()
        # Remove special characters but keep word boundaries
        text_lower = re.sub(r'[^\w\s-]', ' ', text_lower)
        # Normalize whitespace
        text_lower = ' '.join(text_lower.split())
        
        # Count occurrences of keywords for each category
        category_scores = {}
        for category, data in ML_CONCEPTS.items():
            keyword_matches = []
            total_occurrences = 0
            
            for keyword in data['keywords']:
                # Create variations of the keyword
                keyword_variations = [
                    keyword,
                    keyword.replace(' ', '-'),  # hyphenated version
                    keyword.replace(' ', ''),   # no spaces version
                    keyword.replace('-', ' '),  # space instead of hyphen
                ]
                
                # Look for any variation of the keyword
                for variation in keyword_variations:
                    # Look for exact word matches with flexible word boundaries
                    pattern = r'\b' + re.escape(variation.lower()) + r'[s]?\b'
                    matches = re.findall(pattern, text_lower)
                    if matches:
                        if keyword not in keyword_matches:  # Avoid duplicates
                            keyword_matches.append(keyword)
                        total_occurrences += len(matches)
            
            if keyword_matches:
                category_scores[category] = {
                    'matches': keyword_matches,
                    'score': total_occurrences
                }
        
        # If no categories found, check for partial matches
        if not category_scores:
            logger.warning("No exact matches found, checking for partial matches")
            for category, data in ML_CONCEPTS.items():
                for keyword in data['keywords']:
                    # Check if any part of the keyword appears in the text
                    parts = keyword.split()
                    if any(part in text_lower for part in parts if len(part) > 3):
                        if category not in category_scores:
                            category_scores[category] = {
                                'matches': [keyword],
                                'score': 1
                            }
        
        # Build the final concepts dictionary
        for category, score_data in category_scores.items():
            found_keywords = score_data['matches']
            category_data = ML_CONCEPTS[category]
            
            found_concepts[category] = {
                'keywords': found_keywords,
                'definitions': {k: v for k, v in category_data['definitions'].items() if k in found_keywords},
                'relationships': category_data['relationships'],
                'numerical_examples': category_data.get('numerical_examples', {}),
                'practical_cases': category_data.get('practical_cases', {})
            }
        
        # If still no concepts found, return a default concept
        if not found_concepts:
            logger.warning("No ML concepts found, using default concept")
            return {
                'general_ml': {
                    'keywords': ['machine learning'],
                    'definitions': {
                        'machine learning': 'A field of study that gives computers the ability to learn without being explicitly programmed'
                    },
                    'relationships': ['learning', 'prediction'],
                    'numerical_examples': {},
                    'practical_cases': {}
                }
            }
        
        logger.info(f"Found concepts: {list(found_concepts.keys())}")
        return found_concepts
        
    except Exception as e:
        logger.error(f"Error in concept extraction: {str(e)}\n{traceback.format_exc()}")
        return {
            'general_ml': {
                'keywords': ['machine learning'],
                'definitions': {
                    'machine learning': 'A field of study that gives computers the ability to learn without being explicitly programmed'
                },
                'relationships': ['learning', 'prediction'],
                'numerical_examples': {},
                'practical_cases': {}
            }
        }

def generate_numerical_question(concept_data):
    """Generate numerical/calculation based questions"""
    if 'numerical_examples' not in concept_data:
        return None
    
    example = list(concept_data['numerical_examples'].items())[0][1]
    question = {
        'id': f'num-{concept_data["keywords"][0]}',
        'question': f'Consider the following scenario: {example["scenario"]}. Calculate the expected number of samples used for training and explain the final prediction mechanism.',
        'answer_outline': [
            f'Step 1: {example["calculation"]}',
            f'Step 2: {example["result"]}',
            'Step 3: Explain advantages of this approach'
        ],
        'explanation': f'This demonstrates practical application of {concept_data["keywords"][0]} in real scenarios',
        'difficulty': 0.8
    }
    return question

def generate_case_study_question(concept_data):
    """Generate case study based questions"""
    if 'practical_cases' not in concept_data:
        return None
    
    case = list(concept_data['practical_cases'].items())[0][1]
    question = {
        'id': f'case-{concept_data["keywords"][0]}',
        'question': f'''In the following {case["scenario"]}, analyze:
1. Why is {concept_data["keywords"][0]} suitable for this problem?
2. How would you tune the model parameters given current performance: {case["evaluation"]}?
3. What business impact would a false positive/negative have?''',
        'answer_outline': [
            'Problem suitability analysis',
            f'Parameter tuning strategy based on {case["evaluation"]}',
            'Business impact analysis',
            'Recommendations for improvement'
        ],
        'explanation': 'Tests practical understanding and business application',
        'difficulty': 0.9
    }
    return question

def generate_mcq(content, num_questions=4):
    """Generate ML-specific MCQ questions from content"""
    try:
        ml_concepts = extract_ml_concepts(content)
        logger.info(f"MCQ generation - Found concepts: {list(ml_concepts.keys())}")
        
        questions = []
        
        # Template questions to use if we don't have enough concept-based questions
        template_questions = [
            {
                'id': f'mcq-template-1',
                'question': 'Which of the following best describes supervised learning?',
                'options': [
                    'Learning from labeled training data to make predictions',
                    'Finding patterns in unlabeled data',
                    'Learning through trial and error',
                    'Storing data in a database'
                ],
                'correct_answer': 'Learning from labeled training data to make predictions',
                'explanation': 'Supervised learning involves training models on labeled data to make predictions on new, unseen data.',
                'difficulty': 0.7
            },
            {
                'id': f'mcq-template-2',
                'question': 'What is the main purpose of feature selection in machine learning?',
                'options': [
                    'To identify the most relevant input variables for the model',
                    'To increase the number of features',
                    'To make the model more complex',
                    'To slow down training time'
                ],
                'correct_answer': 'To identify the most relevant input variables for the model',
                'explanation': 'Feature selection helps improve model performance by selecting the most informative features.',
                'difficulty': 0.6
            },
            {
                'id': f'mcq-template-3',
                'question': 'What is overfitting in machine learning?',
                'options': [
                    'When a model learns the training data too well and performs poorly on new data',
                    'When a model is too simple to capture patterns',
                    'When the training data is too small',
                    'When the model trains too quickly'
                ],
                'correct_answer': 'When a model learns the training data too well and performs poorly on new data',
                'explanation': 'Overfitting occurs when a model memorizes training data instead of learning general patterns.',
                'difficulty': 0.8
            },
            {
                'id': f'mcq-template-4',
                'question': 'What is the role of the activation function in neural networks?',
                'options': [
                    'To introduce non-linearity into the network',
                    'To speed up training',
                    'To reduce the number of neurons',
                    'To store data'
                ],
                'correct_answer': 'To introduce non-linearity into the network',
                'explanation': 'Activation functions add non-linearity, allowing neural networks to learn complex patterns.',
                'difficulty': 0.7
            }
        ]

        # First try to generate questions from concepts
        for category, data in ml_concepts.items():
            if len(questions) >= num_questions:
                break

            for keyword in data['keywords']:
                if len(questions) >= num_questions:
                    break

                definition = data['definitions'].get(keyword, '')
                if definition:
                    wrong_answers = [
                        f'Process that {data["relationships"][0]} but does not involve {keyword}',
                        f'Technique similar to {data["keywords"][1] if len(data["keywords"]) > 1 else "other methods"} but with different objective',
                        f'Method that only focuses on {data["relationships"][-1]}'
                    ]

                    question = {
                        'id': f'mcq-{len(questions)+1}',
                        'question': f'Which of the following best describes {keyword} in the context of {category.replace("_", " ")}?',
                        'options': [definition] + wrong_answers,
                        'correct_answer': definition,
                        'explanation': f'This is the correct definition because it captures the key aspects of {keyword}.',
                        'difficulty': 0.7
                    }
                    questions.append(question)

        # If we don't have enough questions, add template questions
        while len(questions) < num_questions:
            template_idx = len(questions) % len(template_questions)
            questions.append(template_questions[template_idx])

        return questions[:num_questions]
    except Exception as e:
        logger.error(f"MCQ generation error: {str(e)}\n{traceback.format_exc()}")
        # Return template questions on error
        return template_questions[:num_questions]

def generate_descriptive(content, num_questions=2):
    """Generate ML-specific descriptive questions"""
    try:
        ml_concepts = extract_ml_concepts(content)
        questions = []

        # Template questions to use if needed
        template_questions = [
            {
                'id': 'desc-template-1',
                'question': '''Explain the differences between supervised and unsupervised learning:
1. Define each type of learning
2. Compare their applications
3. Discuss their advantages and limitations
4. Provide real-world examples''',
                'answer_outline': [
                    'Definition of supervised vs unsupervised learning',
                    'Application scenarios for each type',
                    'Pros and cons analysis',
                    'Real-world examples and use cases'
                ],
                'explanation': 'Tests understanding of fundamental ML concepts',
                'difficulty': 0.7
            },
            {
                'id': 'desc-template-2',
                'question': '''Analyze the problem of overfitting in machine learning:
1. What causes overfitting?
2. How can it be detected?
3. What are the methods to prevent it?
4. Discuss the bias-variance tradeoff''',
                'answer_outline': [
                    'Causes of overfitting',
                    'Detection methods',
                    'Prevention techniques',
                    'Bias-variance tradeoff explanation'
                ],
                'explanation': 'Tests understanding of model optimization',
                'difficulty': 0.8
            },
            {
                'id': 'desc-template-3',
                'question': '''Explain the concept of gradient descent:
1. What is the basic principle?
2. Types of gradient descent
3. Role in neural network training
4. Challenges and solutions''',
                'answer_outline': [
                    'Basic principle explanation',
                    'Types: batch, mini-batch, stochastic',
                    'Application in neural networks',
                    'Common challenges and solutions'
                ],
                'explanation': 'Tests understanding of optimization algorithms',
                'difficulty': 0.8
            },
            {
                'id': 'desc-template-4',
                'question': '''Discuss feature selection and engineering:
1. Importance in ML
2. Common techniques
3. Impact on model performance
4. Best practices''',
                'answer_outline': [
                    'Importance of feature selection',
                    'Feature selection techniques',
                    'Performance impact analysis',
                    'Best practices and guidelines'
                ],
                'explanation': 'Tests understanding of data preprocessing',
                'difficulty': 0.7
            }
        ]

        # First try to generate questions from concepts
        for category, data in ml_concepts.items():
            if len(questions) >= num_questions:
                break

            # Add theoretical question
            keywords = data['keywords'][:3]
            if keywords:
                question = {
                    'id': f'desc-{len(questions)+1}',
                    'question': f'''Compare and contrast {", ".join(keywords)} in {category.replace("_", " ")}:
1. Explain the key differences in their approaches
2. Analyze their computational complexity
3. Discuss scenarios where each would be preferred
4. Provide real-world applications''',
                    'answer_outline': [
                        'Theoretical comparison of approaches',
                        'Complexity analysis and performance implications',
                        'Scenario-based selection criteria',
                        'Industry applications and case studies'
                    ],
                    'explanation': f'Tests deep understanding of {category.replace("_", " ")} concepts',
                    'difficulty': 0.8
                }
                questions.append(question)

        # If we don't have enough questions, add template questions
        while len(questions) < num_questions:
            template_idx = len(questions) % len(template_questions)
            questions.append(template_questions[template_idx])

        return questions[:num_questions]
    except Exception as e:
        logger.error(f"Descriptive generation error: {str(e)}")
        return template_questions[:num_questions]

def generate_true_false(content, num_questions=4):
    """Generate ML-specific true/false questions"""
    try:
        ml_concepts = extract_ml_concepts(content)
        questions = []

        # Template questions to use if needed
        template_questions = [
            {
                'id': 'tf-template-1',
                'statement': 'Supervised learning requires labeled training data.',
                'answer': True,
                'explanation': 'Supervised learning algorithms learn from labeled examples to make predictions.',
                'difficulty': 0.6
            },
            {
                'id': 'tf-template-2',
                'statement': 'Neural networks can only be used for classification tasks.',
                'answer': False,
                'explanation': 'Neural networks can be used for various tasks including regression, generation, and more.',
                'difficulty': 0.7
            },
            {
                'id': 'tf-template-3',
                'statement': 'Cross-validation helps prevent overfitting.',
                'answer': True,
                'explanation': 'Cross-validation provides a more robust estimate of model performance.',
                'difficulty': 0.6
            },
            {
                'id': 'tf-template-4',
                'statement': 'Feature scaling is always necessary for all machine learning algorithms.',
                'answer': False,
                'explanation': 'Some algorithms like decision trees are not affected by feature scaling.',
                'difficulty': 0.7
            }
        ]

        # First try to generate questions from concepts
        for category, data in ml_concepts.items():
            if len(questions) >= num_questions:
                break

            for keyword in data['keywords']:
                if len(questions) >= num_questions:
                    break

                definition = data['definitions'].get(keyword, '')
                if definition:
                    question = {
                        'id': f'tf-{len(questions)+1}',
                        'statement': f'{keyword} is primarily used for {data["relationships"][0]}.',
                        'answer': True,
                        'explanation': f'This is a correct statement about the primary use of {keyword}.',
                        'difficulty': 0.6
                    }
                    questions.append(question)

        # If we don't have enough questions, add template questions
        while len(questions) < num_questions:
            template_idx = len(questions) % len(template_questions)
            questions.append(template_questions[template_idx])

        return questions[:num_questions]
    except Exception as e:
        logger.error(f"True/False generation error: {str(e)}")
        return template_questions[:num_questions]

def generate_fill_blank(content, num_questions=4):
    """Generate ML-specific fill in the blank questions"""
    try:
        ml_concepts = extract_ml_concepts(content)
        questions = []

        # Template questions to use if needed
        template_questions = [
            {
                'id': 'fb-template-1',
                'question': '_____ is a type of learning where the model learns from labeled training data.',
                'answer': 'Supervised learning',
                'explanation': 'Supervised learning involves training on labeled data to make predictions.',
                'difficulty': 0.6
            },
            {
                'id': 'fb-template-2',
                'question': 'The _____ function in neural networks introduces non-linearity into the model.',
                'answer': 'activation',
                'explanation': 'Activation functions add non-linearity to neural networks.',
                'difficulty': 0.7
            },
            {
                'id': 'fb-template-3',
                'question': '_____ is the process of finding patterns in unlabeled data.',
                'answer': 'Unsupervised learning',
                'explanation': 'Unsupervised learning discovers patterns without labeled examples.',
                'difficulty': 0.6
            },
            {
                'id': 'fb-template-4',
                'question': 'The _____ algorithm is commonly used to minimize the cost function in neural networks.',
                'answer': 'gradient descent',
                'explanation': 'Gradient descent optimizes the model parameters.',
                'difficulty': 0.7
            }
        ]

        # First try to generate questions from concepts
        for category, data in ml_concepts.items():
            if len(questions) >= num_questions:
                break

            for keyword in data['keywords']:
                if len(questions) >= num_questions:
                    break

                definition = data['definitions'].get(keyword, '')
                if definition:
                    # Replace the keyword with a blank in the definition
                    question_text = definition.replace(keyword, '[BLANK: {}]'.format(keyword))
                    
                    question = {
                        'id': f'fb-{len(questions)+1}',
                        'question': question_text,
                        'answer': keyword,
                        'explanation': f'The term {keyword} fits here as it is being defined.',
                        'difficulty': 0.6,
                        'keyword': keyword
                    }
                    questions.append(question)

        # If we don't have enough questions, add template questions
        while len(questions) < num_questions:
            template_idx = len(questions) % len(template_questions)
            questions.append(template_questions[template_idx])

        return questions[:num_questions]
    except Exception as e:
        logger.error(f"Fill in blank generation error: {str(e)}")
        return template_questions[:num_questions]

def extract_dbms_concepts(text):
    """Extract DBMS-specific concepts from content"""
    dbms_keywords_map = {
        'data_models': {
            'keywords': ['relational model', 'ER model', 'entity-relationship', 'hierarchical model', 'network model', 'object model', 'semi-structured'],
            'definitions': {
                'relational model': 'Uses tables to represent data and relationships',
                'ER model': 'Uses entities and relationships to model data',
                'normalization': 'Process of organizing data to reduce redundancy'
            },
            'relationships': ['data representation', 'schema design', 'modeling']
        },
        'database_design': {
            'keywords': ['normalization', 'functional dependency', 'BCNF', '3NF', '2NF', '1NF', 'normal form', 'decomposition'],
            'definitions': {
                'normalization': 'Process of organizing database tables to reduce redundancy',
                'functional dependency': 'Relationship between attributes in a relation',
                'BCNF': 'Boyce-Codd Normal Form'
            },
            'relationships': ['schema design', 'data integrity', 'redundancy elimination']
        },
        'sql_operations': {
            'keywords': ['SELECT', 'JOIN', 'UNION', 'GROUP BY', 'ORDER BY', 'HAVING', 'WHERE', 'INSERT', 'UPDATE', 'DELETE'],
            'definitions': {
                'JOIN': 'Combines rows from multiple tables',
                'SELECT': 'Retrieves data from database',
                'GROUP BY': 'Groups rows with same values'
            },
            'relationships': ['query processing', 'data retrieval', 'data manipulation']
        },
        'transactions': {
            'keywords': ['ACID', 'transaction', 'atomicity', 'consistency', 'isolation', 'durability', 'concurrency', 'lock', 'deadlock'],
            'definitions': {
                'ACID': 'Properties ensuring reliable transactions',
                'transaction': 'A unit of work performed on database',
                'atomicity': 'All-or-nothing property of transactions'
            },
            'relationships': ['concurrency control', 'data consistency', 'recovery']
        },
        'database_objects': {
            'keywords': ['table', 'view', 'index', 'trigger', 'stored procedure', 'schema', 'constraint', 'primary key', 'foreign key'],
            'definitions': {
                'view': 'Virtual table based on a query',
                'index': 'Structure to speed up data retrieval',
                'constraint': 'Rule enforced on data'
            },
            'relationships': ['database structure', 'data organization', 'performance']
        },
        'database_applications': {
            'keywords': ['OLTP', 'OLAP', 'data warehouse', 'data mining', 'business intelligence', 'ETL'],
            'definitions': {
                'OLTP': 'Online Transaction Processing for day-to-day operations',
                'OLAP': 'Online Analytical Processing for data analysis',
                'data warehouse': 'Central repository for analytical data'
            },
            'relationships': ['data processing', 'analytics', 'decision support']
        }
    }
    
    detected_concepts = {}
    text_lower = text.lower()
    
    for category, data in dbms_keywords_map.items():
        found_keywords = []
        found_definitions = {}
        
        for keyword in data['keywords']:
            if keyword.lower() in text_lower:
                found_keywords.append(keyword)
                if keyword.lower() in data['definitions']:
                    found_definitions[keyword] = data['definitions'][keyword]
        
        if found_keywords:
            detected_concepts[category] = {
                'keywords': found_keywords,
                'definitions': found_definitions,
                'relationships': data['relationships']
            }
    
    return detected_concepts

def generate_summary(text):
    """Generate a comprehensive summary of what the model has learned from the content"""
    try:
        # Detect domain from upload form (if provided)
        course_hint = request.form.get('courseId') if request and request.form else ''
        domain_is_dbms = course_hint and '351A' in course_hint
        domain_is_se = course_hint and '341A' in course_hint

        # Extract concepts based on domain
        if domain_is_dbms:
            domain_concepts = extract_dbms_concepts(text)
            ml_concepts = {}  # No ML concepts for DBMS
        elif domain_is_se:
            domain_concepts = {}  # Could add SE extraction later
            ml_concepts = {}
        else:
            domain_concepts = {}
            ml_concepts = extract_ml_concepts(text)
        
        # Start with an introduction
        summary_sections = []
        
        # First, add any header/title information from the text
        header_lines = []
        for line in text.split('\n')[:10]:  # Look at first 10 lines for headers
            line = line.strip()
            if line and (
                any(word in line.lower() for word in ['department', 'university', 'course', 'subject', 'chapter']) or
                line.isupper() or
                len(line.split()) <= 10  # Short lines are likely headers
            ):
                header_lines.append(line)
        
        if header_lines:
            summary_sections.extend(header_lines)
            summary_sections.append("")  # Add spacing
        
        # Add document statistics
        summary_sections.append("# Document Analysis")
        
        # Split into paragraphs and clean up
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        words = text.split()
        sentences = [s.strip() for s in text.split('.') if s.strip()]
        
        summary_sections.append(f"- Total Paragraphs: {len(paragraphs)}")
        summary_sections.append(f"- Total Words: {len(words)}")
        summary_sections.append(f"- Total Sentences: {len(sentences)}")
        summary_sections.append(f"- Average Words per Sentence: {len(words)/max(1, len(sentences)):.1f}")
        summary_sections.append("")
        
        # Add main concepts overview - use domain_concepts for DBMS/SE, ml_concepts for ML
        concepts_to_show = domain_concepts if domain_is_dbms or domain_is_se else ml_concepts
        
        if concepts_to_show:
            # Course-aware heading
            if domain_is_se:
                summary_sections.append("# Detected SE Concepts")
            elif domain_is_dbms:
                summary_sections.append("# Detected DBMS Concepts")
            else:
                summary_sections.append("# Detected ML Concepts")
            
            for category, data in concepts_to_show.items():
                category_name = category.replace('_', ' ').title()
                found_keywords = data['keywords']
                
                summary_sections.append(f"\n## {category_name}")
                
                # Add key concepts and their definitions
                summary_sections.append("### Key Concepts Identified:")
                for keyword in found_keywords:
                    if keyword in data.get('definitions', {}):
                        summary_sections.append(f"- **{keyword}**: {data['definitions'][keyword]}")
                    else:
                        summary_sections.append(f"- **{keyword}**")
                
                # Add relationships between concepts
                if len(found_keywords) > 1 and 'relationships' in data:
                    summary_sections.append("\n### Concept Relationships:")
                    for rel in data['relationships']:
                        summary_sections.append(f"- {rel.title()}")
                
                # Add any practical applications if available (ML only)
                if 'practical_cases' in data and data['practical_cases'] and not domain_is_dbms and not domain_is_se:
                    summary_sections.append("\n### Practical Applications:")
                    for case_name, case in data['practical_cases'].items():
                        summary_sections.append(f"- {case['scenario']}")
                        if 'evaluation' in case:
                            metrics = [f"{k}: {v}" for k, v in case['evaluation'].items()]
                            summary_sections.append(f"  - Performance Metrics: {', '.join(metrics)}")
                
                summary_sections.append("")  # Add spacing between categories
        else:
            # If nothing detected, keep neutral wording
            if domain_is_dbms:
                summary_sections.append("\n# Detected DBMS Concepts")
                summary_sections.append("No specific DBMS concepts were detected in the content. Please ensure the slides contain DBMS-related terminology.")
            elif domain_is_se:
                summary_sections.append("\n# Detected SE Concepts")
                summary_sections.append("No specific SE concepts were detected in the content.")
            else:
                summary_sections.append("\n# Note: Concepts Not Detected")
                summary_sections.append("The content did not match the expected ML concepts.")
        
        # Add content analysis
        summary_sections.append("# Content Analysis")
        
        # Analyze complexity level
        concepts_for_complexity = concepts_to_show if concepts_to_show else {}
        technical_terms = sum(len(data['keywords']) for data in concepts_for_complexity.values()) if concepts_for_complexity else 0
        if technical_terms > 20:
            complexity = "Advanced"
        elif technical_terms > 10:
            complexity = "Intermediate"
        else:
            complexity = "Basic"
        summary_sections.append(f"- Content Complexity: {complexity}")
        
        # Analyze coverage (domain-aware wording)
        if domain_is_dbms:
            if concepts_to_show:
                covered_categories = len(concepts_to_show)
                summary_sections.append(f"- Topic Coverage: {covered_categories} DBMS concept categories detected")
            else:
                summary_sections.append("- Topic Coverage: Based on DBMS keywords extracted from slides")
        elif domain_is_se:
            if concepts_to_show:
                covered_categories = len(concepts_to_show)
                summary_sections.append(f"- Topic Coverage: {covered_categories} SE concept categories detected")
            else:
                summary_sections.append("- Topic Coverage: Based on SE keywords extracted from slides")
        else:
            if ml_concepts:
                covered_categories = len(ml_concepts)
                total_categories = len(ML_CONCEPTS)
                coverage = (covered_categories / total_categories) * 100
                summary_sections.append(f"- Topic Coverage: {coverage:.1f}% of core ML concepts")
            else:
                summary_sections.append("- Topic Coverage: Based on ML keywords extracted from slides")
        
        # Extract potential prerequisites (ML only)
        prerequisites = set()
        if not domain_is_dbms and not domain_is_se:
            for category, data in ml_concepts.items():
                if category in ['neural_networks', 'supervised_learning']:
                    prerequisites.add('Linear Algebra')
                    prerequisites.add('Calculus')
                if category in ['model_optimization', 'evaluation_metrics']:
                    prerequisites.add('Statistics')
                    prerequisites.add('Probability')
        
        if prerequisites:
            summary_sections.append("\n## Prerequisites")
            for prereq in sorted(prerequisites):
                summary_sections.append(f"- {prereq}")
        
        # Add recommendations for question generation
        if concepts_to_show:
            summary_sections.append("\n# Question Generation Strategy")
            summary_sections.append("Based on the content analysis, questions will focus on:")
            for category in concepts_to_show:
                category_name = category.replace('_', ' ').title()
                summary_sections.append(f"- {category_name}")
        
        # Add a note about the current context
        summary_sections.append("\n# Processing Information")
        summary_sections.append(f"- Timestamp: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        summary_sections.append(f"- Content Hash: {hash(text) % 10000:04d}")
        
        # Add the full text at the end
        summary_sections.append("\n# Full Content")
        summary_sections.append(text)
        
        return '\n'.join(summary_sections)
    except Exception as e:
        logger.error(f"Summary generation error: {str(e)}")
        return "Error generating summary"

def generate_questions_without_peft(content, question_type, num_questions):
    """Generate questions using the traditional methods as fallback"""
    if question_type == 'mcq':
        return generate_mcq(content, num_questions)
    elif question_type == 'descriptive':
        return generate_descriptive(content, num_questions)
    elif question_type == 'true_false':
        return generate_true_false(content, num_questions)
    else:  # fill_blank
        return generate_fill_blank(content, num_questions)

def generate_questions_with_peft(content, question_type, num_questions):
    """Generate questions using the PEFT-adapted model"""
    try:
        # Extract ML concepts from content
        ml_concepts = extract_ml_concepts(content)
        logger.info(f"Found concepts for question generation: {list(ml_concepts.keys())}")
        
        # Get the topics from the image for prioritization
        priority_topics = [
            'cnn', 'bagging', 'decision_trees', 'ensemble_learning',
            'map', 'svm', 'hmm', 'em', 'unsupervised'
        ]
        
        # Filter and sort concepts by priority
        available_concepts = []
        
        # First add concepts that match the priority topics
        for topic in priority_topics:
            normalized_topic = topic.replace(' ', '_').lower()
            for concept_key in ml_concepts.keys():
                if normalized_topic in concept_key.lower():
                    available_concepts.append((concept_key, ml_concepts[concept_key]))
                    break
        
        # Then add any remaining detected concepts
        for topic, data in ml_concepts.items():
            if not any(topic.lower() in [t.replace(' ', '_').lower() for t, _ in available_concepts]):
                available_concepts.append((topic, data))
        
        if not available_concepts:
            logger.warning("No specific ML concepts found, falling back to traditional generation")
            return generate_questions_without_peft(content, question_type, num_questions)
        
        questions = []
        concepts_used = set()
        
        # Generate questions ensuring topic diversity
        while len(questions) < num_questions and available_concepts:
            # Rotate through available concepts
            for topic, concept_data in available_concepts:
                if len(questions) >= num_questions:
                    break
                    
                if topic in concepts_used and len(available_concepts) > 1:
                    continue
                
                logger.info(f"Generating {question_type} question for topic: {topic}")
                
                # Get template based on question type and topic
                if question_type == 'mcq':
                    question = generate_topic_specific_mcq(topic, concept_data)
                elif question_type == 'descriptive':
                    question = generate_topic_specific_descriptive(topic, concept_data)
                elif question_type == 'true_false':
                    question = generate_topic_specific_true_false(topic, concept_data)
                else:  # fill_blank
                    question = generate_topic_specific_fill_blank(topic, concept_data)
                
                if question:
                    questions.append(question)
                    concepts_used.add(topic)
            
            # If we still need more questions, allow reusing topics
            if len(questions) < num_questions:
                concepts_used.clear()
        
        return questions
    except Exception as e:
        logger.error(f"Error in PEFT question generation: {str(e)}")
        return generate_questions_without_peft(content, question_type, num_questions)

def generate_topic_specific_mcq(topic, concept_data):
    """Generate high-quality MCQ with detailed technical content using Bloom's Taxonomy"""
    try:
        # Get a concept and its definition
        concept = random.choice(concept_data['keywords'])
        definition = concept_data['definitions'].get(concept)
        
        # If we don't have a proper definition, try another concept
        if not definition or len(definition) < 20:
            for alt_concept in concept_data['keywords']:
                if alt_concept in concept_data['definitions']:
                    definition = concept_data['definitions'][alt_concept]
                    concept = alt_concept
                    if len(definition) >= 20:
                        break
        
        # If still no good definition, return None
        if not definition or len(definition) < 20:
            return None

        # Select a random Bloom's level, heavily weighted towards lower levels (easy questions)
        bloom_weights = [0.5, 0.3, 0.1, 0.05, 0.03, 0.02]  # Favor remember/understand (easy)
        bloom_level = random.choices(list(BLOOMS_TAXONOMY.keys()), weights=bloom_weights)[0]
        taxonomy_data = BLOOMS_TAXONOMY[bloom_level]

        # Determine difficulty based on Bloom's level - more conservative
        if taxonomy_data['level'] <= 2:  # remember, understand
            difficulty = 'easy'
        elif taxonomy_data['level'] <= 3:  # apply only
            difficulty = 'moderate'
        else:  # analyze, evaluate, create
            difficulty = 'hard'

        # Generate question using Bloom's template
        question_template = random.choice(taxonomy_data['mcq_templates'])
        question_text = question_template.format(f"'{concept}'")

        # Create distractors based on Bloom's level
        distractors = []
        
        if bloom_level == 'remember':
            # Basic factual distractors
            if len(concept_data['keywords']) > 1:
                related_concept = random.choice([k for k in concept_data['keywords'] if k != concept])
                distractors.append(concept_data['definitions'].get(related_concept, f"A different type of {topic.replace('_', ' ')}"))
            distractors.append(f"A process unrelated to {topic.replace('_', ' ')}")
            distractors.append(f"A technique that does not involve {concept}")
            
        elif bloom_level == 'understand':
            # Understanding-based distractors
            distractors.append(f"{definition.split(' but ')[0]} but with different purpose")
            distractors.append(f"Similar to {concept} but used in different context")
            distractors.append(f"A misconception about how {concept} works")
            
        elif bloom_level == 'apply':
            # Application-based distractors
            distractors.append(f"An incorrect implementation of {concept}")
            distractors.append(f"A suboptimal way to apply {concept}")
            distractors.append(f"A common mistake when using {concept}")
            
        elif bloom_level == 'analyze':
            # Analysis-based distractors
            distractors.append(f"An incomplete analysis of {concept}'s behavior")
            distractors.append(f"A superficial examination of {concept}'s effects")
            distractors.append(f"A misinterpretation of {concept}'s role")
            
        elif bloom_level == 'evaluate':
            # Evaluation-based distractors
            distractors.append(f"An incorrect evaluation of {concept}'s effectiveness")
            distractors.append(f"A biased assessment of {concept}'s benefits")
            distractors.append(f"A flawed comparison of {concept} with alternatives")
            
        else:  # create
            # Creation-based distractors
            distractors.append(f"An impractical design using {concept}")
            distractors.append(f"An inefficient implementation of {concept}")
            distractors.append(f"A flawed innovation involving {concept}")

        # Shuffle options but remember correct answer index
        all_options = [definition] + distractors[:3]
        correct_index = 0
        random.shuffle(all_options)
        correct_index = all_options.index(definition)
        
        return {
            'id': f'mcq-{topic}-{hash(concept) % 1000:03d}',
            'question': question_text,
            'keyword': concept,
            'options': all_options,
            'correct_answer': definition,
            'answer': correct_index,
            'explanation': f"This question tests {bloom_level}-level understanding of {concept}",
            'difficulty': difficulty,
            'topic': topic,
            'concept': concept,
            'bloom_level': bloom_level,
            'cognitive_demand': taxonomy_data['level']
        }

    except Exception as e:
        logger.error(f"Error generating MCQ: {str(e)}")
        return None

def generate_topic_specific_descriptive(topic, concept_data):
    """Generate high-quality descriptive questions with technical depth using Bloom's Taxonomy"""
    try:
        # Select 2-3 related concepts for comparison
        concepts = random.sample(concept_data['keywords'], min(3, len(concept_data['keywords'])))
        
        # Select Bloom's level - descriptive can be moderate but still favor easier
        bloom_weights = [0.2, 0.3, 0.25, 0.15, 0.07, 0.03]  # Still favor understand/apply
        bloom_level = random.choices(list(BLOOMS_TAXONOMY.keys()), weights=bloom_weights)[0]
        taxonomy_data = BLOOMS_TAXONOMY[bloom_level]

        # Determine difficulty - more conservative
        if taxonomy_data['level'] <= 2:  # remember, understand
            difficulty = 'easy'
        elif taxonomy_data['level'] <= 3:  # apply
            difficulty = 'moderate'
        else:  # analyze, evaluate, create
            difficulty = 'hard'

        # Create question text and outline
        question_template = random.choice(taxonomy_data['descriptive_templates'])
        question_text = question_template.format(
            concepts[0],
            concepts[1] if len(concepts) > 1 else topic.replace('_', ' ')
        )

        # Build comprehensive answer using concept definitions
        concept_def = concept_data['definitions'].get(concepts[0], '')
        related_concepts = [c for c in concepts[1:] if c in concept_data['definitions']]
        relationships = concept_data.get('relationships', [])
        
        answer_parts = []
        
        # Part 1: Definition
        if concept_def:
            answer_parts.append(f"Definition: {concept_def}")
        else:
            answer_parts.append(f"Definition: {concepts[0]} is a key concept in {topic.replace('_', ' ')}")
        
        # Part 2: Key components
        if related_concepts:
            answer_parts.append(f"Key Components: {concepts[0]} involves {', '.join(related_concepts[:2])}")
        if relationships:
            answer_parts.append(f"Relationships: {concepts[0]} is related to {', '.join(relationships[:2])}")
        
        # Part 3: Examples
        if 'practical_cases' in concept_data and concept_data['practical_cases']:
            case = list(concept_data['practical_cases'].values())[0]
            answer_parts.append(f"Example: {case.get('scenario', 'Practical application of the concept')}")
        else:
            answer_parts.append(f"Example: {concepts[0]} can be applied in real-world scenarios")
        
        # Part 4: Applications
        answer_parts.append(f"Applications: {concepts[0]} is used for {relationships[0] if relationships else 'various applications'} in {topic.replace('_', ' ')}")
        
        answer_text = ' '.join(answer_parts)
        answer_outline = [
            f"Define and explain {concepts[0]}",
            "Analyze key components and relationships",
            "Provide examples and applications",
            "Evaluate effectiveness and limitations"
        ]
        
        return {
            'id': f'desc-{topic}-{hash("".join(concepts)) % 1000:03d}',
            'question': question_text,
            'keyword': concepts[0],
            'answer': answer_text,
            'answer_outline': answer_outline,
            'correct_answer': answer_text,
            'explanation': f"This question tests {bloom_level}-level understanding. The answer should comprehensively cover the definition, components, examples, and applications of {concepts[0]}.",
            'difficulty': difficulty,
            'topic': topic,
            'concept': concepts[0],
            'bloom_level': bloom_level,
            'cognitive_demand': taxonomy_data['level']
        }

    except Exception as e:
        logger.error(f"Error generating descriptive question: {str(e)}")
        return None

def generate_topic_specific_true_false(topic, concept_data):
    """Generate true/false questions with technical accuracy"""
    try:
        # Get a concept and its relationships
        concept = random.choice(concept_data['keywords'])
        relationships = concept_data.get('relationships', [])
        definition = concept_data['definitions'].get(concept, '')

        if not definition:
            return None

        # Select a random Bloom's level, weighted towards easy
        bloom_weights = [0.5, 0.3, 0.1, 0.05, 0.03, 0.02]  # Favor remember/understand
        bloom_level = random.choices(list(BLOOMS_TAXONOMY.keys()), weights=bloom_weights)[0]
        taxonomy_data = BLOOMS_TAXONOMY[bloom_level]

        # Determine difficulty - more conservative
        if taxonomy_data['level'] <= 2:  # remember, understand
            difficulty = 'easy'
        elif taxonomy_data['level'] <= 3:  # apply
            difficulty = 'moderate'
        else:  # analyze, evaluate, create
            difficulty = 'hard'

        # Create true statement
        true_statement = f"{concept} is used for {relationships[0] if relationships else 'processing'} in {topic.replace('_', ' ')}."
        
        # Create false statement by modifying the true statement
        false_statement = f"{concept} completely eliminates the need for {relationships[-1] if relationships else 'other techniques'} in {topic.replace('_', ' ')}."

        # Randomly choose true or false statement
        is_true = random.choice([True, False])
        statement = true_statement if is_true else false_statement

        # Build comprehensive explanation with correct answer
        if is_true:
            explanation = f"This statement is TRUE. {definition}. {concept} is indeed used for {relationships[0] if relationships else 'processing'} in {topic.replace('_', ' ')}."
        else:
            explanation = f"This statement is FALSE. {concept} does not completely eliminate the need for {relationships[-1] if relationships else 'other techniques'}. {definition}."

        return {
            'id': f'tf-{topic}-{hash(statement) % 1000:03d}',
            'question': statement,
            'statement': statement,
            'answer': is_true,
            'correct_answer': is_true,
            'correct_answer_text': 'True' if is_true else 'False',
            'explanation': explanation,
            'difficulty': difficulty,
            'topic': topic,
            'concept': concept,
            'bloom_level': bloom_level
        }

    except Exception as e:
        logger.error(f"Error generating true/false question: {str(e)}")
        return None

def generate_topic_specific_fill_blank(topic, concept_data):
    """Generate high-quality fill-in-the-blank questions with technical context"""
    try:
        # Get concept and technical details
        concept = random.choice(concept_data['keywords'])
        definition = concept_data['definitions'].get(concept, '')
        relationships = concept_data.get('relationships', [])

        if not definition:
            return None

        # Select a random Bloom's level, weighted towards easy
        bloom_weights = [0.5, 0.3, 0.1, 0.05, 0.03, 0.02]  # Favor remember/understand
        bloom_level = random.choices(list(BLOOMS_TAXONOMY.keys()), weights=bloom_weights)[0]
        taxonomy_data = BLOOMS_TAXONOMY[bloom_level]

        # Determine difficulty - more conservative
        if taxonomy_data['level'] <= 2:  # remember, understand
            difficulty = 'easy'
        elif taxonomy_data['level'] <= 3:  # apply
            difficulty = 'moderate'
        else:  # analyze, evaluate, create
            difficulty = 'hard'

        # Create question text with blank
        question_text = f"In {topic.replace('_', ' ')}, [BLANK] is used for {relationships[0] if relationships else 'processing'} because {definition.split('.')[0]}."

        # Build comprehensive explanation
        explanation = f"The correct answer is '{concept}'. {definition}. In {topic.replace('_', ' ')}, {concept} is used for {relationships[0] if relationships else 'processing'}."

        return {
            'id': f'fb-{topic}-{hash(question_text) % 1000:03d}',
            'question': question_text,
            'answer': concept,
            'correct_answer': concept,
            'keyword': concept,
            'explanation': explanation,
            'difficulty': difficulty,
            'topic': topic,
            'concept': concept,
            'bloom_level': bloom_level
        }

    except Exception as e:
        logger.error(f"Error generating fill-in-blank question: {str(e)}")
        return None

def generate_questions_from_summary(summary_text, question_counts):
    """Generate questions directly from the summary text"""
    try:
        logger.info("Generating questions from summary...")
        logger.info(f"Requested question counts (max 4 per type): {question_counts}")
        
        # Extract topic and subtopics from summary
        lines = summary_text.split('\n')
        main_topic = lines[0] if lines else ""
        subtopics = [line.strip() for line in lines[1:] if line.strip()]
        
        logger.info(f"Main topic: {main_topic}")
        logger.info(f"Subtopics: {subtopics}")
        
        # Find matching topic in ML_CONCEPTS with more flexible matching
        detected_topics = []
        for topic, data in ML_CONCEPTS.items():
            # Check main topic
            if any(keyword.lower() in main_topic.lower() for keyword in data['keywords']):
                detected_topics.append((topic, 1.0))  # Primary match
                continue
            
            # Check subtopics
            for subtopic in subtopics:
                if any(keyword.lower() in subtopic.lower() for keyword in data['keywords']):
                    detected_topics.append((topic, 0.8))  # Secondary match
                    break
        
        if not detected_topics:
            logger.warning("No specific topic detected in summary")
            return None
        
        # Sort topics by match confidence and use the best match
        detected_topics.sort(key=lambda x: x[1], reverse=True)
        detected_topic, confidence = detected_topics[0]
            
        logger.info(f"Detected primary topic: {detected_topic} (confidence: {confidence})")
        concept_data = ML_CONCEPTS[detected_topic]
        
        questions = {
                'MCQs': [],
                'Descriptive': [],
                'TrueFalse': [],
            'FillBlank': [],
            'CaseStudy': [],
            'Numerical': []
        }
        
        # Generate MCQs (max 4)
        mcq_count = min(4, question_counts.get('mcq', 4))
        for _ in range(mcq_count):
            mcq = generate_topic_specific_mcq(detected_topic, concept_data)
            if mcq:
                questions['MCQs'].append(mcq)
        
        # Generate descriptive questions (max 4)
        desc_count = min(4, question_counts.get('descriptive', 4))
        for _ in range(desc_count):
            desc = generate_topic_specific_descriptive(detected_topic, concept_data)
            if desc:
                questions['Descriptive'].append(desc)
        
        # Generate true/false questions (max 4)
        tf_count = min(4, question_counts.get('true_false', 4))
        for _ in range(tf_count):
            tf = generate_topic_specific_true_false(detected_topic, concept_data)
            if tf:
                questions['TrueFalse'].append(tf)
        
        # Generate fill in blank questions (max 4)
        fb_count = min(4, question_counts.get('fill_in_blanks', 4))
        for _ in range(fb_count):
            fb = generate_topic_specific_fill_blank(detected_topic, concept_data)
            if fb:
                questions['FillBlank'].append(fb)
        
        # Generate case study questions if available (max 4)
        if 'practical_cases' in concept_data:
            case_count = min(4, question_counts.get('case_study', 1))
            for _ in range(case_count):
                case = generate_case_study_question(detected_topic, concept_data)
                if case:
                    questions['CaseStudy'].append(case)
        
        # Generate numerical questions if available (max 4)
        if 'numerical_examples' in concept_data:
            num_count = min(4, question_counts.get('numerical', 1))
            for _ in range(num_count):
                num = generate_numerical_question(detected_topic, concept_data)
                if num:
                    questions['Numerical'].append(num)
        
        # Remove empty question types
        return {k: v for k, v in questions.items() if v}
        
    except Exception as e:
        logger.error(f"Error generating questions from summary: {str(e)}")
        return None

def generate_dbms_question_from_content(content, question_type, domain='DBMS'):
    """Generate domain-specific questions from content using improved prompt engineering"""
    try:
        import re
        import random
        
        # Extract key DBMS/SE terms from content
        dbms_keywords = ['ACID', 'normalization', 'SQL', 'JOIN', 'transaction', 'database', 
                        'relational', 'schema', 'index', 'query', 'constraint', 'ER model',
                        'OLTP', 'OLAP', 'data model', 'table', 'row', 'column', 'primary key',
                        'foreign key', 'trigger', 'view', 'stored procedure', 'deadlock',
                        'concurrency', 'isolation level', 'two-phase locking', 'Entity-Relationship',
                        'normal form', 'BCNF', 'functional dependency', 'integrity']
        
        se_keywords = ['requirement', 'UML', 'design', 'testing', 'unit test', 'integration',
                      'system', 'software', 'architecture', 'pattern', 'agile', 'scrum',
                      'validation', 'verification', 'SDLC', 'waterfall', 'iterative']
        
        keywords = dbms_keywords if domain == 'DBMS' else se_keywords
        
        # Find which keywords appear in content with context
        content_lower = content.lower()
        found_keywords = []
        for kw in keywords:
            if kw.lower() in content_lower:
                # Extract sentence containing the keyword for context
                sentences = [s.strip() for s in content.split('.') if kw.lower() in s.lower()]
                if sentences:
                    found_keywords.append((kw, sentences[0]))
        
        # If no keywords found, extract noun phrases from content
        if not found_keywords:
            words = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', content)
            found_keywords = [(w, '') for w in list(set([w for w in words if len(w) > 3]))[:10]]
        
        if not found_keywords:
            found_keywords = [('database concepts', '')]
        
        # Select a random topic from found keywords (avoid repetition)
        if not hasattr(generate_dbms_question_from_content, 'used_topics'):
            generate_dbms_question_from_content.used_topics = []
        
        available_topics = [t for t in found_keywords if t[0] not in generate_dbms_question_from_content.used_topics]
        if not available_topics:
            generate_dbms_question_from_content.used_topics = []  # Reset
            available_topics = found_keywords
        
        topic_tuple = random.choice(available_topics) if available_topics else ('database', '')
        topic = topic_tuple[0]
        topic_context = topic_tuple[1]
        generate_dbms_question_from_content.used_topics.append(topic)
        
        # Determine difficulty based on topic complexity
        def get_difficulty(topic_str):
            topic_l = topic_str.lower()
            # Easy: basic definitions, simple concepts
            if any(x in topic_l for x in ['table', 'row', 'column', 'database', 'data', 'sql']):
                return 'easy'
            # Moderate: intermediate concepts
            elif any(x in topic_l for x in ['join', 'view', 'index', 'constraint', 'schema', 'primary key', 'foreign key']):
                return 'moderate'
            # Hard: advanced concepts
            elif any(x in topic_l for x in ['normalization', 'acid', 'transaction', 'concurrency', 'deadlock', 'isolation', 'two-phase']):
                return 'moderate'  # Still moderate for most
            else:
                return 'easy'
        
        difficulty = get_difficulty(topic)
        
        # Extract sentences mentioning the topic for all question types
        topic_sentences = [s.strip() for s in content.split('.') 
                         if topic.lower() in s.lower() and len(s.strip()) > 20]
        
        # Generate question based on type
        if question_type == 'mcq':
            # Create MCQ with domain-specific options and correct answers
            if domain == 'DBMS':
                if 'acid' in topic.lower():
                    # Enhance with content if available
                    acid_content = ' '.join([s for s in topic_sentences if any(x in s.lower() for x in ['atomic', 'consist', 'isolat', 'durabl'])][:2])
                    explanation_base = 'Atomicity ensures that either all operations in a transaction complete successfully, or none do. If any part fails, the entire transaction is rolled back.'
                    if acid_content:
                        explanation_base = f"{explanation_base} {acid_content[:150]}"
                    
                    # Create question about ACID properties
                    acid_questions = [
                        {
                            'question': 'Which ACID property ensures that a transaction is treated as a single, indivisible unit?',
                            'options': ['Atomicity', 'Consistency', 'Isolation', 'Durability'],
                            'correct_answer': 'Atomicity',
                            'answer': 0,
                            'explanation': explanation_base,
                            'difficulty': 'easy'
                        },
                        {
                            'question': 'Which ACID property ensures that database constraints are maintained before and after transactions?',
                            'options': ['Atomicity', 'Consistency', 'Isolation', 'Durability'],
                            'correct_answer': 'Consistency',
                            'answer': 1,
                            'explanation': 'Consistency ensures that the database remains in a valid state, following all defined rules and constraints, both before and after the transaction.',
                            'difficulty': 'easy'
                        },
                        {
                            'question': 'Which ACID property prevents transactions from interfering with each other?',
                            'options': ['Atomicity', 'Consistency', 'Isolation', 'Durability'],
                            'correct_answer': 'Isolation',
                            'answer': 2,
                            'explanation': 'Isolation ensures that concurrent transactions execute independently without interfering with each other, as if they were running sequentially.',
                            'difficulty': 'easy'
                        },
                        {
                            'question': 'Which ACID property ensures that committed changes persist even after system failures?',
                            'options': ['Atomicity', 'Consistency', 'Isolation', 'Durability'],
                            'correct_answer': 'Durability',
                            'answer': 3,
                            'explanation': 'Durability guarantees that once a transaction is committed, its changes are permanently saved and will survive system crashes or failures.',
                            'difficulty': 'easy'
                        }
                    ]
                    return random.choice(acid_questions)
                    
                elif 'normal' in topic.lower() or 'normal form' in topic.lower():
                    return {
                        'question': 'What is the primary goal of database normalization?',
                        'options': [
                            'To eliminate data redundancy and improve data integrity',
                            'To increase database size for faster access',
                            'To speed up all queries without any trade-offs',
                            'To remove all foreign keys from the database'
                        ],
                        'correct_answer': 'To eliminate data redundancy and improve data integrity',
                        'answer': 0,
                        'explanation': 'Normalization reduces data redundancy by organizing data into related tables, which prevents update anomalies and improves data integrity.',
                        'difficulty': 'easy'
                    }
                elif 'join' in topic.lower() or 'sql' in topic.lower():
                    sql_questions = [
                        {
                            'question': 'Which SQL operation combines rows from two or more tables based on a related column?',
                            'options': ['GROUP BY', 'JOIN', 'ORDER BY', 'UNION'],
                            'correct_answer': 'JOIN',
                            'answer': 1,
                            'explanation': 'JOIN combines rows from multiple tables based on a common column or relationship, allowing you to retrieve related data from different tables.',
                            'difficulty': 'easy'
                        },
                        {
                            'question': 'What does SQL stand for?',
                            'options': [
                                'Structured Query Language',
                                'Simple Query Logic',
                                'System Query Language',
                                'Sequential Query Language'
                            ],
                            'correct_answer': 'Structured Query Language',
                            'answer': 0,
                            'explanation': 'SQL (Structured Query Language) is a standardized language used to manage and query relational databases.',
                            'difficulty': 'easy'
                        },
                        {
                            'question': 'What is the difference between INNER JOIN and OUTER JOIN?',
                            'options': [
                                'INNER JOIN returns only matching rows, OUTER JOIN returns all rows from one or both tables',
                                'INNER JOIN is faster than OUTER JOIN',
                                'OUTER JOIN is used only for text data',
                                'There is no difference'
                            ],
                            'correct_answer': 'INNER JOIN returns only matching rows, OUTER JOIN returns all rows from one or both tables',
                            'answer': 0,
                            'explanation': 'INNER JOIN returns only rows where there is a match in both tables, while OUTER JOIN (LEFT, RIGHT, FULL) includes unmatched rows from one or both tables.',
                            'difficulty': 'moderate'
                        },
                        {
                            'question': 'Which SQL clause is used to filter rows after grouping?',
                            'options': ['WHERE', 'HAVING', 'GROUP BY', 'ORDER BY'],
                            'correct_answer': 'HAVING',
                            'answer': 1,
                            'explanation': 'HAVING is used to filter groups created by GROUP BY, while WHERE filters rows before grouping.',
                            'difficulty': 'easy'
                        }
                    ]
                    return random.choice(sql_questions)
                    
                elif 'transaction' in topic.lower():
                    return {
                        'question': 'What does ACID stand for in database transactions?',
                        'options': [
                            'Atomicity, Consistency, Isolation, Durability',
                            'Access, Control, Integrity, Data',
                            'Analysis, Calculation, Index, Database',
                            'Application, Command, Interface, Design'
                        ],
                        'correct_answer': 'Atomicity, Consistency, Isolation, Durability',
                        'answer': 0,
                        'explanation': 'ACID is an acronym for the four key properties that guarantee reliable transaction processing: Atomicity, Consistency, Isolation, and Durability.',
                        'difficulty': 'easy'
                    }
            
            # Extract actual content about the topic
            topic_sentences = [s.strip() for s in content.split('.') 
                             if topic.lower() in s.lower() and len(s.strip()) > 20]
            
            if topic_sentences:
                # Use actual content to create question
                context_sentence = topic_sentences[0][:150]  # First relevant sentence
                # Create a question based on the actual content
                if 'relational' in topic.lower() or 'model' in topic.lower():
                    return {
                        'question': 'Which of the following best describes the relational data model?',
                        'options': [
                            'A collection of tables representing data and relationships',
                            'A hierarchical tree structure',
                            'A network of nodes and edges',
                            'A flat file structure'
                        ],
                        'correct_answer': 'A collection of tables representing data and relationships',
                        'answer': 0,
                        'explanation': 'The relational model uses tables (relations) to represent both data and relationships.',
                        'difficulty': 'easy'
                    }
                elif 'oltp' in topic.lower() or 'olap' in topic.lower():
                    return {
                        'question': 'What is the main difference between OLTP and OLAP?',
                        'options': [
                            'OLTP handles transactions, OLAP handles analytics',
                            'OLTP is faster than OLAP in all cases',
                            'OLAP uses only SQL, OLTP uses NoSQL',
                            'There is no functional difference'
                        ],
                        'correct_answer': 'OLTP handles transactions, OLAP handles analytics',
                        'answer': 0,
                        'explanation': 'OLTP (Online Transaction Processing) handles day-to-day operational transactions, while OLAP (Online Analytical Processing) focuses on data analysis and decision support.',
                        'difficulty': 'easy'
                    }
                elif 'er' in topic.lower() or 'entity' in topic.lower():
                    return {
                        'question': 'What does the Entity-Relationship model represent?',
                        'options': [
                            'Entities and relationships between them',
                            'Only database tables',
                            'Only SQL queries',
                            'Network topology'
                        ],
                        'correct_answer': 'Entities and relationships between them',
                        'answer': 0,
                        'explanation': 'The ER model uses entities (real-world objects) and relationships to represent the conceptual structure of data.',
                        'difficulty': 'easy'
                    }
                elif 'data model' in topic.lower():
                    return {
                        'question': 'Which data model uses tables to represent data and relationships?',
                        'options': ['Relational model', 'Hierarchical model', 'Network model', 'Object model'],
                        'correct_answer': 'Relational model',
                        'answer': 0,
                        'explanation': 'The relational model uses tables (relations) with rows (tuples) and columns (attributes) to represent data and relationships.',
                        'difficulty': 'easy'
                    }
                elif 'schema' in topic.lower():
                    return {
                        'question': 'What is a database schema?',
                        'options': [
                            'The structure and organization of the database',
                            'A type of query',
                            'A backup file',
                            'A user interface'
                        ],
                        'correct_answer': 'The structure and organization of the database',
                        'answer': 0,
                        'explanation': 'A schema defines the structure, tables, relationships, constraints, and other elements of a database.',
                        'difficulty': 'easy'
                    }
                elif 'view' in topic.lower():
                    return {
                        'question': 'What is a database view?',
                        'options': [
                            'A virtual table based on a query',
                            'A physical table stored separately',
                            'A type of index',
                            'A backup mechanism'
                        ],
                        'correct_answer': 'A virtual table based on a query',
                        'answer': 0,
                        'explanation': 'A view is a virtual table that presents data from underlying tables based on a stored query, without storing the data itself.',
                        'difficulty': 'easy'
                    }
                elif 'constraint' in topic.lower():
                    constraint_questions = [
                        {
                            'question': 'What is the purpose of database constraints?',
                            'options': [
                                'To enforce data integrity and rules',
                                'To speed up all queries automatically',
                                'To reduce storage space',
                                'To create automatic backups'
                            ],
                            'correct_answer': 'To enforce data integrity and rules',
                            'answer': 0,
                            'explanation': 'Constraints ensure data follows rules like primary keys (uniqueness), foreign keys (referential integrity), and check constraints (value validation).',
                            'difficulty': 'easy'
                        },
                        {
                            'question': 'What is the difference between a primary key and a foreign key?',
                            'options': [
                                'Primary key uniquely identifies a row in a table, foreign key references primary key in another table',
                                'Primary key is always numeric, foreign key is always text',
                                'Primary key is optional, foreign key is mandatory',
                                'There is no difference'
                            ],
                            'correct_answer': 'Primary key uniquely identifies a row in a table, foreign key references primary key in another table',
                            'answer': 0,
                            'explanation': 'A primary key uniquely identifies each row in its own table, while a foreign key in one table references the primary key of another table to establish relationships.',
                            'difficulty': 'easy'
                        }
                    ]
                    return random.choice(constraint_questions)
                    
                elif 'primary key' in topic.lower() or 'foreign key' in topic.lower():
                    return {
                        'question': 'What is the main purpose of a primary key in a database table?',
                        'options': [
                            'To uniquely identify each row in the table',
                            'To speed up all queries',
                            'To store encrypted data',
                            'To create automatic indexes'
                        ],
                        'correct_answer': 'To uniquely identify each row in the table',
                        'answer': 0,
                        'explanation': 'A primary key ensures each row has a unique identifier, which is essential for relational database integrity and relationships.',
                        'difficulty': 'easy'
                    }
                    
                elif 'index' in topic.lower():
                    return {
                        'question': 'What is the primary purpose of a database index?',
                        'options': [
                            'To speed up data retrieval operations',
                            'To reduce storage space',
                            'To encrypt data',
                            'To create backups'
                        ],
                        'correct_answer': 'To speed up data retrieval operations',
                        'answer': 0,
                        'explanation': 'Indexes create a data structure that allows faster lookup of rows based on indexed columns, significantly improving query performance.',
                        'difficulty': 'easy'
                    }
                    
                elif 'concurrency' in topic.lower() or 'lock' in topic.lower():
                    return {
                        'question': 'What problem does concurrency control solve in databases?',
                        'options': [
                            'Prevents conflicts when multiple transactions access the same data simultaneously',
                            'Reduces database size',
                            'Speeds up all queries',
                            'Creates automatic backups'
                        ],
                        'correct_answer': 'Prevents conflicts when multiple transactions access the same data simultaneously',
                        'answer': 0,
                        'explanation': 'Concurrency control ensures that when multiple transactions run simultaneously, they do not interfere with each other, maintaining data consistency.',
                        'difficulty': 'moderate'
                    }
            
            # Content-aware MCQ using actual slide content
            if topic_sentences:
                # Extract key information from content
                content_snippet = topic_sentences[0][:150]
                # Clean leading bullet points from content snippet
                content_snippet = content_snippet.lstrip('•- ').strip()
                # Create question based on actual content - use exam-style format
                if 'used' in content_snippet.lower() or 'purpose' in content_snippet.lower():
                    question = f'What is the primary purpose of {topic} in database systems?'
                    # Extract purpose from content
                    correct_option = content_snippet.split('.')[0] if '.' in content_snippet else content_snippet
                    # Clean leading bullet points from correct option
                    correct_option = correct_option.lstrip('•- ').strip()
                else:
                    question = f'Which of the following best describes {topic}?'
                    correct_option = f'{topic} is a key concept in {domain} for {content_snippet[:80]}'
                    # Clean leading bullet points from correct option
                    correct_option = correct_option.lstrip('•- ').strip()
            else:
                question = f'What is the primary purpose of {topic} in {domain}?'
                correct_option = f'Managing and organizing data in {domain} systems'
            
            # Clean bullet points from all options
            options = [
                correct_option,
                f'Hardware configuration in computer systems',
                f'Network routing and protocol management',
                f'User interface design and development'
            ]
            # Clean leading bullet points from options (preserve dashes within text)
            options = [opt.lstrip('•- ').strip() if isinstance(opt, str) else opt for opt in options]
            
            return {
                'question': question,
                'options': options,
                'correct_answer': options[0],  # Use cleaned option
                'answer': 0,
                'explanation': f'{topic} is a fundamental concept in {domain}. {topic_sentences[0][:100] if topic_sentences else "It is used for managing and organizing data."}',
                'difficulty': difficulty
            }
        
        elif question_type == 'descriptive':
            # Extract actual content about the topic for better questions
            topic_sentences = [s.strip() for s in content.split('.') 
                             if topic.lower() in s.lower() and len(s.strip()) > 30]
            
            # Create domain-specific answer templates based on actual content
            if domain == 'DBMS':
                if 'acid' in topic.lower():
                    # Extract ACID-related content
                    acid_content = ' '.join([s for s in topic_sentences if any(x in s.lower() for x in ['atomic', 'consist', 'isolat', 'durabl'])][:3])
                    if acid_content:
                        answer_text = f"ACID properties ensure reliable database transactions: (1) Atomicity - all operations succeed or all fail together, (2) Consistency - database remains valid before and after transaction, (3) Isolation - concurrent transactions don't interfere, (4) Durability - committed changes persist even after failures. {acid_content[:200]}"
                    else:
                        answer_text = f"ACID properties ensure reliable database transactions: (1) Atomicity - all operations succeed or all fail together, (2) Consistency - database remains valid before and after transaction, (3) Isolation - concurrent transactions don't interfere, (4) Durability - committed changes persist even after failures. Example: When transferring money between accounts, both debit and credit must complete together (atomicity), and the transaction must be recorded permanently (durability)."
                elif 'normal' in topic.lower():
                    norm_content = ' '.join([s for s in topic_sentences if 'normal' in s.lower() or 'redund' in s.lower()][:2])
                    if norm_content:
                        answer_text = f"Normalization is the process of organizing database tables to reduce redundancy: (1) Eliminates data duplication, (2) Prevents update anomalies (insertion, deletion, modification), (3) Ensures data integrity, (4) Normal forms (1NF, 2NF, 3NF, BCNF) define levels of normalization. {norm_content[:200]}"
                    else:
                        answer_text = f"Normalization is the process of organizing database tables to reduce redundancy: (1) Eliminates data duplication, (2) Prevents update anomalies (insertion, deletion, modification), (3) Ensures data integrity, (4) Normal forms (1NF, 2NF, 3NF, BCNF) define levels of normalization. Example: Instead of storing customer address in every order record, create a separate Customers table and reference it via foreign key."
                elif 'join' in topic.lower():
                    join_content = ' '.join([s for s in topic_sentences if 'join' in s.lower()][:2])
                    if join_content:
                        answer_text = f"SQL JOIN combines rows from multiple tables: (1) INNER JOIN returns matching rows from both tables, (2) LEFT/RIGHT JOIN returns all rows from one table and matches from other, (3) FULL JOIN returns all rows from both tables, (4) JOINs are based on common columns (usually primary/foreign keys). {join_content[:200]}"
                    else:
                        answer_text = f"SQL JOIN combines rows from multiple tables: (1) INNER JOIN returns matching rows from both tables, (2) LEFT/RIGHT JOIN returns all rows from one table and matches from other, (3) FULL JOIN returns all rows from both tables, (4) JOINs are based on common columns (usually primary/foreign keys). Example: JOIN Customers and Orders tables on customer_id to get customer names with their orders."
                elif 'transaction' in topic.lower():
                    trans_content = ' '.join([s for s in topic_sentences if 'transaction' in s.lower()][:2])
                    if trans_content:
                        answer_text = f"A transaction is a unit of work that must execute completely or not at all: (1) Groups multiple database operations into one atomic operation, (2) Ensures data consistency even if system fails, (3) Uses ACID properties. {trans_content[:200]}"
                    else:
                        answer_text = f"A transaction is a unit of work that must execute completely or not at all: (1) Groups multiple database operations into one atomic operation, (2) Ensures data consistency even if system fails, (3) Uses ACID properties. Example: Transferring money involves debiting one account and crediting another - both must succeed together."
                else:
                    # Use actual content about the topic
                    if topic_sentences:
                        content_excerpt = topic_sentences[0][:250]
                        answer_text = f"{topic} in DBMS: (1) Definition - {content_excerpt}, (2) Purpose - Used for managing and organizing database operations, (3) Key characteristics and relationships, (4) Practical applications in database systems."
                    elif topic_context:
                        answer_text = f"{topic} in DBMS: (1) Definition - {topic_context[:200]}, (2) Purpose - Used for managing and organizing database operations, (3) Key characteristics, (4) Practical applications."
                    else:
                        answer_text = f"{topic} in DBMS: (1) Definition and purpose, (2) Key characteristics, (3) How it relates to database operations, (4) Practical applications in database systems."
            else:
                # ML descriptive - use actual content
                if topic_sentences:
                    content_excerpt = topic_sentences[0][:250]
                    answer_text = f"{topic} in {domain}: (1) Definition - {content_excerpt}, (2) Key characteristics and purpose, (3) How it works and its relationships, (4) Applications and use cases."
                else:
                    answer_text = f"{topic} in {domain}: (1) Definition and purpose, (2) Key characteristics, (3) How it works, (4) Applications and use cases."
            
            # Create question based on actual content - use exam-style format
            if topic_sentences:
                question_text = f"Explain the concept of {topic} in {domain}. Provide a detailed explanation with examples."
            else:
                question_text = f"Explain the concept of {topic} in {domain}. Provide detailed explanation with examples."
            
            return {
                'question': question_text,
                'answer': answer_text,
                'correct_answer': answer_text,
                'answer_outline': [
                    f'Definition of {topic}',
                    'Key characteristics and purpose',
                    'How it works and relationships',
                    'Applications and use cases'
                ],
                'explanation': f'This question requires understanding of {topic}. The answer should cover the definition, characteristics, how it works, and applications.',
                'difficulty': difficulty
            }
        
        elif question_type == 'true_false':
            # Create domain-specific T/F statements with correct answers
            if domain == 'DBMS':
                if 'acid' in topic.lower():
                    statements = [
                        ('ACID properties ensure reliable transaction processing in databases.', True),
                        ('Atomicity guarantees that a transaction is treated as a single, indivisible unit.', True),
                        ('ACID stands for Association, Consistency, Isolation, and Durability.', False),
                        ('Durability ensures that committed transactions can be rolled back.', False)
                    ]
                elif 'normal' in topic.lower():
                    statements = [
                        ('Normalization eliminates data redundancy in database tables.', True),
                        ('The goal of normalization is to increase database size.', False),
                        ('BCNF (Boyce-Codd Normal Form) is a higher normal form than 3NF.', True),
                        ('Normalization always improves query performance.', False)
                    ]
                elif 'join' in topic.lower():
                    statements = [
                        ('JOIN combines rows from multiple tables based on related columns.', True),
                        ('INNER JOIN returns only matching rows from both tables.', True),
                        ('OUTER JOIN is faster than INNER JOIN in all cases.', False),
                        ('JOIN operations require primary keys in both tables.', False)
                    ]
                elif 'transaction' in topic.lower():
                    statements = [
                        ('A transaction must follow ACID properties.', True),
                        ('Transactions ensure data consistency even if the system fails.', True),
                        ('All database operations must be part of a transaction.', False),
                        ('Transactions can be partially committed.', False)
                    ]
                else:
                    # Use content-based statements if available
                    if topic_sentences:
                        statements = [
                            (f'{topic} is an important concept in {domain}.', True),
                            (f'{topic_sentences[0][:100]}...', True) if len(topic_sentences[0]) > 50 else (f'{topic} is used in {domain} systems.', True)
                        ]
                    else:
                        statements = [
                            (f'{topic} is an important concept in {domain}.', True),
                            (f'{topic} is used for managing data in {domain} systems.', True)
                        ]
            else:
                # ML T/F - use content if available
                if topic_sentences:
                    statements = [
                        (f'{topic} is an important concept in {domain}.', True),
                        (f'{topic_sentences[0][:100]}...', True) if len(topic_sentences[0]) > 50 else (f'{topic} is used in {domain}.', True)
                    ]
                else:
                    statements = [
                        (f'{topic} is an important concept in {domain}.', True),
                        (f'{topic} is used in {domain} systems.', True)
                    ]
            
            statement, correct_answer = random.choice(statements)
            
            # Create explanation using content if available
            if topic_sentences:
                explanation = f'This statement is {str(correct_answer).lower()}. {topic_sentences[0][:150]}'
            else:
                explanation = f'This statement is {str(correct_answer).lower()}. {topic} is a key concept in {domain}.'
            
            return {
                'question': statement,
                'statement': statement,
                'answer': correct_answer,
                'correct_answer': correct_answer,
                'correct_answer_text': 'True' if correct_answer else 'False',
                'explanation': explanation,
                'difficulty': difficulty
            }
        
        elif question_type == 'fill_blank':
            # Create better fill-in-the-blank based on topic with comprehensive explanations
            if 'acid' in topic.lower():
                question = f'In database transactions, the ACID property of ____ ensures all-or-nothing execution.'
                answer = 'Atomicity'
                explanation = f"The correct answer is 'Atomicity'. Atomicity ensures that a transaction is treated as a single, indivisible unit - either all operations complete successfully, or none do. This is one of the four ACID properties that guarantee reliable transaction processing."
            elif 'normal' in topic.lower():
                question = f'Database ____ reduces data redundancy by organizing data into related tables.'
                answer = 'normalization'
                explanation = f"The correct answer is 'normalization'. Normalization is the process of organizing database tables to eliminate data redundancy by breaking tables into smaller, related tables and using foreign keys to maintain relationships."
            elif 'join' in topic.lower():
                question = f'The SQL ____ operation combines rows from multiple tables based on a common column.'
                answer = 'JOIN'
                explanation = f"The correct answer is 'JOIN'. JOIN is a SQL operation that combines rows from two or more tables based on a related column (typically a primary key-foreign key relationship), allowing retrieval of related data from multiple tables."
            else:
                question = f'____ is a key concept in {domain} related to {topic}.'
                answer = topic
                if topic_sentences:
                    explanation = f"The correct answer is '{answer}'. {topic_sentences[0][:150]}"
                else:
                    explanation = f"The correct answer is '{answer}'. {topic} is a fundamental concept in {domain} used for managing and organizing data."
            
            return {
                'question': question,
                'answer': answer,
                'correct_answer': answer,
                'explanation': explanation,
                'difficulty': difficulty
            }
        
        return None
    except Exception as e:
        logger.error(f"Error generating {domain} question: {str(e)}")
        return None

@app.route('/generate-questions', methods=['POST'])
def generate_questions_endpoint():
    try:
        start_time = time.time()
        data = request.get_json()
        
        # Load content from request or file
        content = data.get('content')
        if not content:
            content = load_context()
        
        if not content:
            logger.error("No content available for question generation")
            return jsonify({
                'success': False, 
                'error': 'No content available. Please upload content first.'
            }), 400

        # Log content for debugging
        logger.info(f"Content length: {len(content)}")
        logger.info(f"Content preview: {content[:200]}...")

        # Extract concepts (subject-aware)
        course_hint = (data.get('courseId') or '').upper()
        is_dbms = '351A' in course_hint
        is_se = '341A' in course_hint
        
        # For DBMS/SE, use prompt-based generation from content instead of ML concept extraction
        if is_dbms or is_se:
            # Use content directly with domain-specific prompts
            concepts = None  # Will use direct content generation
        else:
            concepts = extract_ml_concepts(content)

        # Get requested counts
        requested_counts = data.get('counts', {
            'mcq': 4,
            'descriptive': 4,
            'true_false': 4,
            'fill_in_blanks': 4
        })

        # Initialize questions
        questions = {
            'MCQs': [],
            'Descriptive': [],
            'TrueFalse': [],
            'FillBlank': []
        }

        # For DBMS/SE, use prompt-based generation from actual content
        if is_dbms or is_se:
            logger.info(f"Using prompt-based generation for {'DBMS' if is_dbms else 'SE'}")
            domain = 'DBMS' if is_dbms else 'SE'
            
            # Use full content for better question generation
            content_sample = content
            
            # Reset topic tracker for this request
            generate_dbms_question_from_content.used_topics = []
            
            # Generate MCQs using prompt engineering
            mcq_count = requested_counts.get('mcq', 4)
            for i in range(mcq_count):
                mcq = generate_dbms_question_from_content(content_sample, 'mcq', domain)
                if mcq:
                    # Add ID if not present
                    if 'id' not in mcq:
                        mcq['id'] = f'mcq-{len(questions["MCQs"]) + 1}'
                    # Clean bullet points from options (only leading bullets)
                    if 'options' in mcq and isinstance(mcq['options'], list):
                        mcq['options'] = [opt.lstrip('•- ').strip() if isinstance(opt, str) else opt for opt in mcq['options']]
                    questions['MCQs'].append(mcq)
            
            # Generate Descriptive questions
            desc_count = requested_counts.get('descriptive', 4)
            for i in range(desc_count):
                desc = generate_dbms_question_from_content(content_sample, 'descriptive', domain)
                if desc:
                    # Add ID if not present
                    if 'id' not in desc:
                        desc['id'] = f'desc-{len(questions["Descriptive"]) + 1}'
                    questions['Descriptive'].append(desc)
            
            # Generate True/False questions
            tf_count = requested_counts.get('true_false', 4)
            for i in range(tf_count):
                tf = generate_dbms_question_from_content(content_sample, 'true_false', domain)
                if tf:
                    # Add ID if not present
                    if 'id' not in tf:
                        tf['id'] = f'tf-{len(questions["TrueFalse"]) + 1}'
                    questions['TrueFalse'].append(tf)
            
            # Generate Fill in Blank questions
            fb_count = requested_counts.get('fill_in_blanks', 4)
            for i in range(fb_count):
                fb = generate_dbms_question_from_content(content_sample, 'fill_blank', domain)
                if fb:
                    # Add ID if not present
                    if 'id' not in fb:
                        fb['id'] = f'fb-{len(questions["FillBlank"]) + 1}'
                    questions['FillBlank'].append(fb)
            
            concepts_list = ['Content-based topics']
        else:
            # ML path - use existing concept extraction
            logger.info(f"Found ML concepts: {list(concepts.keys())}")
            
            # Generate questions for each concept
            for concept_name, concept_data in concepts.items():
                logger.info(f"Generating questions for concept: {concept_name}")
                
                # Generate MCQs
                for _ in range(requested_counts.get('mcq', 4)):
                    mcq = generate_topic_specific_mcq(concept_name, concept_data)
                    if mcq and isinstance(mcq, dict):
                        questions['MCQs'].append(mcq)
                        logger.info(f"Generated MCQ for {concept_name}")

                # Generate descriptive questions
                for _ in range(requested_counts.get('descriptive', 4)):
                    desc = generate_topic_specific_descriptive(concept_name, concept_data)
                    if desc and isinstance(desc, dict):
                        questions['Descriptive'].append(desc)
                        logger.info(f"Generated descriptive question for {concept_name}")

                # Generate true/false questions
                for _ in range(requested_counts.get('true_false', 4)):
                    tf = generate_topic_specific_true_false(concept_name, concept_data)
                    if tf and isinstance(tf, dict):
                        questions['TrueFalse'].append(tf)
                        logger.info(f"Generated true/false question for {concept_name}")

                # Generate fill in blank questions
                for _ in range(requested_counts.get('fill_in_blanks', 4)):
                    fb = generate_topic_specific_fill_blank(concept_name, concept_data)
                    if fb and isinstance(fb, dict):
                        questions['FillBlank'].append(fb)
                        logger.info(f"Generated fill in blank question for {concept_name}")
            
            concepts_list = list(concepts.keys())

        # Log generation results
        total_questions = sum(len(q) for q in questions.values())
        logger.info(f"Generated {total_questions} questions total")
        for qtype, qlist in questions.items():
            logger.info(f"Generated {len(qlist)} {qtype}")

        if total_questions == 0:
            logger.error("No questions were generated")
            return jsonify({
                'success': False,
                'error': 'Failed to generate questions. Please try uploading different content or check if the content contains sufficient machine learning concepts.'
            }), 400

        # Calculate statistics
        stats = {
            'total_questions': total_questions,
            'type_distribution': {
                qtype: len(qlist) for qtype, qlist in questions.items()
            },
            'concepts_covered': concepts_list if 'concepts_list' in locals() else (list(concepts.keys()) if concepts else [])
        }

        return jsonify({
            'success': True,
            'questions': questions,
            'stats': stats,
            'concepts_found': concepts_list if 'concepts_list' in locals() else (list(concepts.keys()) if concepts else []),
            'message': f"Successfully generated {total_questions} questions!"
        })

    except Exception as e:
        logger.error(f"Question generation failed: {str(e)}\n{traceback.format_exc()}")
        return jsonify({
            'success': False, 
            'error': f'Error generating questions: {str(e)}',
            'suggestions': [
                'Make sure the uploaded content contains machine learning related text',
                'Try uploading a longer document with more detailed content',
                'Check if the file format is supported (PDF, TXT, DOCX, or PPTX)'
            ]
        }), 500

@app.route('/generate-paper', methods=['POST'])
def generate_paper():
    try:
        data = request.get_json()
        logger.info(f"Received paper generation request: {data}")
        
        # Validate request data
        required_fields = ['title', 'totalMarks', 'duration', 'questionTypes']
        if not all(field in data for field in required_fields):
            return jsonify({
                'success': False,
                'error': 'Missing required fields'
            }), 400

        # Create papers directory if it doesn't exist
        os.makedirs(app.config['PAPERS_DIR'], exist_ok=True)

        # Map frontend question types to database types (lowercase)
        type_mapping = {
            'mcq': 'mcq',
            'descriptive': 'descriptive',
            'trueFalse': 'true_false',
            'fillInBlanks': 'fill_blanks'
        }
        
        # Get questions from MongoDB based on selected types
        selected_types = [type_mapping[qtype] for qtype, selected in data['questionTypes'].items() 
                         if selected and qtype in type_mapping]
        
        logger.info(f"Selected question types: {selected_types}")
        
        # Connect to MongoDB
        client = MongoClient('mongodb://localhost:27017/')
        db = client['pesuprep']
        questions_collection = db['questiontemplates']
        
        # Get random questions of selected types
        questions = []
        marks_per_question = data['totalMarks'] // 10  # Assuming we want 10 questions
        
        for qtype in selected_types:
            logger.info(f"Fetching questions of type: {qtype}")
            type_questions = list(questions_collection.find({
                'type': qtype,
                'status': 'approved'
            }).limit(3))
            
            logger.info(f"Found {len(type_questions)} questions of type {qtype}")
            if type_questions:
                questions.extend(type_questions)
        
        if not questions:
            logger.error("No questions found in database")
            return jsonify({
                'success': False,
                'error': 'No questions available for selected types'
            }), 400

        # Randomly select up to 10 questions
        selected_questions = random.sample(questions, min(10, len(questions)))
        logger.info(f"Selected {len(selected_questions)} questions for the paper")

        # Create paper document
        paper = {
            'title': data['title'],
            'description': data.get('description', ''),
            'instructions': data.get('instructions', ''),
            'totalMarks': data['totalMarks'],
            'duration': data['duration'],
            'questions': selected_questions,
            'created_at': datetime.now().isoformat(),
            'status': 'active'
        }

        # Save paper to papers collection
        papers_collection = db['papers']
        result = papers_collection.insert_one(paper)
        paper_id = str(result.inserted_id)
        logger.info(f"Paper saved with ID: {paper_id}")

        # Generate PDF
        pdf_path = os.path.join(app.config['PAPERS_DIR'], f'{paper_id}.pdf')
        
        # Create PDF
        generate_paper_pdf(paper, pdf_path)
        logger.info(f"PDF generated at: {pdf_path}")

        return jsonify({
            'success': True,
            'paper_id': paper_id,
            'message': 'Paper generated successfully'
        })

    except Exception as e:
        logger.error(f"Error generating paper: {str(e)}\n{traceback.format_exc()}")
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/papers/<paper_id>', methods=['GET'])
def get_paper(paper_id):
    try:
        if not paper_id or paper_id == 'undefined':
            return jsonify({
                'success': False,
                'error': 'Invalid paper ID'
            }), 400
            
        pdf_path = os.path.join(app.config['PAPERS_DIR'], f'{paper_id}.pdf')
        
        if not os.path.exists(pdf_path):
            logger.error(f"PDF file not found: {pdf_path}")
            return jsonify({
                'success': False,
                'error': 'Paper not found'
            }), 404

        return send_file(pdf_path, mimetype='application/pdf')
    except Exception as e:
        logger.error(f"Error retrieving paper: {str(e)}")
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/papers/<paper_id>', methods=['DELETE'])
def delete_paper(paper_id):
    try:
        if not paper_id or paper_id == 'undefined':
            return jsonify({
                'success': False,
                'error': 'Invalid paper ID'
            }), 400
            
        # Connect to MongoDB
        client = MongoClient('mongodb://localhost:27017/')
        db = client['pesuprep']
        papers_collection = db['papers']
        
        # Delete from database
        result = papers_collection.delete_one({'_id': ObjectId(paper_id)})
        
        if result.deleted_count == 0:
            return jsonify({
                'success': False,
                'error': 'Paper not found'
            }), 404

        # Delete PDF file
        pdf_path = os.path.join(app.config['PAPERS_DIR'], f'{paper_id}.pdf')
        if os.path.exists(pdf_path):
            os.remove(pdf_path)

        return jsonify({
            'success': True,
            'message': 'Paper deleted successfully'
        })

    except Exception as e:
        logger.error(f"Error deleting paper: {str(e)}")
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

def generate_paper_pdf(paper, output_path):
    """Generate PDF for the question paper"""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    
    # Create the document with margins
    doc = SimpleDocTemplate(
        output_path,
        pagesize=letter,
        leftMargin=0.75*inch,
        rightMargin=0.75*inch,
        topMargin=0.75*inch,
        bottomMargin=0.75*inch
    )
    
    styles = getSampleStyleSheet()
    
    # Create custom styles
    styles.add(ParagraphStyle(
        name='CustomTitle',
        parent=styles['Title'],
        fontSize=16,
        spaceAfter=20,
        alignment=1,  # Center alignment
        textColor=colors.black
    ))
    
    styles.add(ParagraphStyle(
        name='CustomHeading',
        parent=styles['Heading1'],
        fontSize=12,
        spaceAfter=12,
        textColor=colors.black
    ))
    
    styles.add(ParagraphStyle(
        name='Instructions',
        parent=styles['Normal'],
        fontSize=10,
        spaceAfter=12,
        textColor=colors.black
    ))
    
    # Build document content
    content = []
    
    # Add title
    content.append(Paragraph(paper['title'], styles['CustomTitle']))
    
    # Add metadata in a better formatted table
    metadata = [
        ['Duration:', f"{paper['duration']} minutes"],
        ['Total Marks:', str(paper['totalMarks'])]
    ]
    
    meta_table = Table(metadata, colWidths=[1.5*inch, 4*inch])
    meta_table.setStyle(TableStyle([
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('TEXTCOLOR', (0, 0), (-1, -1), colors.black),
        ('FONTSIZE', (0, 0), (-1, -1), 11),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
        ('TOPPADDING', (0, 0), (-1, -1), 8),
        ('FONT', (0, 0), (0, -1), 'Helvetica-Bold'),
        ('FONT', (1, 0), (1, -1), 'Helvetica'),
    ]))
    content.append(meta_table)
    content.append(Spacer(1, 20))
    
    # Add instructions
    if paper.get('instructions'):
        content.append(Paragraph('Instructions:', styles['CustomHeading']))
        content.append(Paragraph(paper['instructions'], styles['Instructions']))
        content.append(Spacer(1, 20))
    
    # Add questions
    for i, question in enumerate(paper['questions'], 1):
        # Question number and text
        content.append(Paragraph(
            f"Q{i}. {question['questionText']}",
            styles['CustomHeading']
        ))
        
        # Add options for MCQ
        if question.get('options'):
            for j, option in enumerate(question['options']):
                option_text = f"{chr(65+j)}. {option['text']}"
                content.append(Paragraph(option_text, styles['Normal']))
        
        content.append(Spacer(1, 15))
    
    # Generate PDF
    doc.build(content)

def serialize_mongo_doc(doc):
    """Helper function to serialize MongoDB document"""
    if isinstance(doc, dict):
        return {k: serialize_mongo_doc(v) for k, v in doc.items()}
    elif isinstance(doc, list):
        return [serialize_mongo_doc(item) for item in doc]
    elif isinstance(doc, ObjectId):
        return str(doc)
    elif isinstance(doc, datetime):
        return doc.isoformat()
    return doc

@app.route('/papers', methods=['GET'])
def get_papers():
    try:
        # Connect to MongoDB
        client = MongoClient('mongodb://localhost:27017/')
        db = client['pesuprep']
        papers_collection = db['papers']
        
        # Get all papers, sorted by creation date
        papers = list(papers_collection.find({}).sort('created_at', -1))
        
        # Serialize the papers to make them JSON-friendly
        serialized_papers = [serialize_mongo_doc(paper) for paper in papers]
        
        return jsonify({
            'success': True,
            'papers': serialized_papers
        })
    except Exception as e:
        logger.error(f"Error fetching papers: {str(e)}")
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

def get_bloom_level(question_text):
    """Determine the Bloom's taxonomy level of a question based on keywords"""
    question_lower = question_text.lower()
    max_level = 0
    assigned_level = 'remember'  # default level
    
    for level, data in BLOOMS_TAXONOMY.items():
        if any(keyword in question_lower for keyword in data['keywords']):
            if data['level'] > max_level:
                max_level = data['level']
                assigned_level = level
    
    return assigned_level

def adapt_question_difficulty(question_data, student_level):
    """Adapt question difficulty based on student's cognitive level"""
    try:
        if not isinstance(question_data, dict):
            logger.error(f"Invalid question data type: {type(question_data)}")
            return None

        # Extract question text
        question_text = question_data.get('question', '')
        if not question_text:
            logger.error("Question text not found in question data")
            return None

        # Get current Bloom's level
        current_level = get_bloom_level(question_text)
        if not current_level:
            logger.error("Could not determine current Bloom's level")
            return question_data  # Return original question if level can't be determined

        # Convert levels to numeric values (1-6)
        level_map = {
            'remember': 1,
            'understand': 2,
            'apply': 3,
            'analyze': 4,
            'evaluate': 5,
            'create': 6
        }

        current_numeric = level_map[current_level]
        target_numeric = min(max(1, int(student_level * 6)), 6)  # Convert 0-1 to 1-6 range

        # If levels match or difference is small, return original
        if abs(current_numeric - target_numeric) <= 1:
            return question_data

        # Get target level name
        target_level = next(key for key, value in level_map.items() if value == target_numeric)

        # Generate new question at target level
        topic = question_data.get('topic', '')
        concept = question_data.get('concept', '')
        question_type = 'mcq' if 'options' in question_data else 'descriptive'

        if not topic or not concept:
            logger.error("Missing topic or concept in question data")
            return question_data

        # Create concept data structure
        concept_data = {
            'topic': topic,
            'keywords': [concept],
            'definitions': {concept: question_data.get('explanation', '')}
        }

        # Generate new question at target level
        new_question = generate_bloom_question(topic, concept, target_level, question_type)
        
        if new_question:
            # Preserve metadata from original question
            new_question['id'] = question_data.get('id', '')
            new_question['topic'] = topic
            new_question['concept'] = concept
            return new_question

        return question_data  # Fallback to original if adaptation fails

    except Exception as e:
        logger.error(f"Error adapting question difficulty: {str(e)}")
        return question_data  # Return original question on error

def get_bloom_level(question_text):
    """Determine the Bloom's taxonomy level of a question based on keywords"""
    # Define keyword patterns for each level
    bloom_patterns = {
        'remember': r'\b(define|list|recall|name|identify|show|what|when|where|who|describe)\b',
        'understand': r'\b(explain|interpret|summarize|infer|compare|classify|how|why)\b',
        'apply': r'\b(implement|solve|use|demonstrate|apply|calculate|complete|illustrate)\b',
        'analyze': r'\b(analyze|differentiate|organize|relate|compare|contrast|examine|experiment)\b',
        'evaluate': r'\b(evaluate|assess|critique|judge|justify|defend|determine)\b',
        'create': r'\b(create|design|construct|develop|compose|plan|propose|formulate)\b'
    }

    # Count matches for each level
    level_counts = {}
    for level, pattern in bloom_patterns.items():
        matches = len(re.findall(pattern, question_text.lower()))
        level_counts[level] = matches

    # Return the level with the most keyword matches
    if any(level_counts.values()):
        return max(level_counts.items(), key=lambda x: x[1])[0]
    
    # Default to 'understand' if no keywords found
    return 'understand'

def generate_bloom_question(topic, concept, level, question_type='mcq'):
    """Generate a question at a specific Bloom's taxonomy level"""
    try:
        # Create basic concept data
        concept_data = {
            'topic': topic,
            'keywords': [concept],
            'definitions': {},
            'relationships': []
        }

        # Generate question based on type and level
        if question_type == 'mcq':
            return generate_topic_specific_mcq(topic, concept_data)
        elif question_type == 'descriptive':
            return generate_topic_specific_descriptive(topic, concept_data)
        elif question_type == 'true_false':
            return generate_topic_specific_true_false(topic, concept_data)
        elif question_type == 'fill_blank':
            return generate_topic_specific_fill_blank(topic, concept_data)
        else:
            logger.error(f"Unsupported question type: {question_type}")
            return None

    except Exception as e:
        logger.error(f"Error generating Bloom's question: {str(e)}")
        return None

# Add signal handlers for tracking downtime
import signal

def signal_handler(signum, frame):
    """Track system downtime"""
    METRICS['downtime_periods'].append((datetime.now(), datetime.now() + timedelta(seconds=1)))
    
signal.signal(signal.SIGTERM, signal_handler)
signal.signal(signal.SIGINT, signal_handler)

@app.route('/api/questions/review', methods=['POST'])
def review_question():
    try:
        data = request.get_json()
        question_id = data.get('questionId')
        status = data.get('status')
        feedback = data.get('feedback')
        timestamp = data.get('timestamp')
        
        if not question_id:
            return jsonify({'message': 'Question ID is required'}), 400
            
        # Track teacher validation
        METRICS['teacher_validations']['total'] += 1
        if status == 'approved':
            METRICS['teacher_validations']['approved'] += 1

        # Connect to MongoDB
        client = MongoClient('mongodb://localhost:27017/')
        db = client['pesuprep']
        question_templates = db['questiontemplates']

        # Update the question's status in questiontemplates collection
        result = question_templates.find_one_and_update(
            {'_id': ObjectId(question_id)},
            {
                '$set': {
                    'status': status,
                    'feedback': feedback,
                    'reviewedAt': timestamp or datetime.now()
                }
            },
            return_document=True  # Return the updated document
        )

        if not result:
            return jsonify({ 
                'message': 'Question not found',
                'question': { 'id': question_id, 'status': status },
                'metrics': calculate_metrics()
            }), 404

        # Convert ObjectId to string for JSON serialization
        result['_id'] = str(result['_id'])

        return jsonify({
            'message': 'Question reviewed successfully',
            'question': result,
            'metrics': calculate_metrics()
        })

    except Exception as e:
        logger.error(f"Review error: {str(e)}")
        return jsonify({
            'message': 'Error reviewing question',
            'error': str(e)
        }), 500

def filterByDifficulty(questions, difficultyLevel):
    """Filter questions by difficulty level using string categories"""
    return [q for q in questions if q.get('difficulty', 'moderate') == difficultyLevel]

@app.route('/evaluate-answers', methods=['POST'])
def evaluate_answers():
    try:
        data = request.get_json()
        if not data or 'answers' not in data:
            return jsonify({'error': 'No answers provided'}), 400

        answers = data['answers']
        results = []

        for answer in answers:
            student_answer = answer.get('answer', '')
            correct_answer = answer.get('correctAnswer', '')
            question_type = answer.get('type', '')
            confidence_score = answer.get('confidenceScore', 1)

            # For MCQ and true/false, evaluate directly
            if question_type in ['mcq', 'trueFalse']:
                is_correct = student_answer == correct_answer
                results.append({
                    'questionId': answer.get('questionId', ''),
                    'evaluation': 'Correct' if is_correct else 'Incorrect',
                    'score': 100 if is_correct else 0,
                    'isCorrect': is_correct,
                    'confidenceScore': confidence_score
                })
                continue

            # For fill in blanks, do simple string comparison
            if question_type == 'fillInBlanks':
                # Case-insensitive comparison
                is_correct = student_answer.lower().strip() == correct_answer.lower().strip()
                results.append({
                    'questionId': answer.get('questionId', ''),
                    'evaluation': 'Correct' if is_correct else 'Incorrect',
                    'score': 100 if is_correct else 0,
                    'isCorrect': is_correct,
                    'confidenceScore': confidence_score
                })
                continue

            # For descriptive questions, use simpler evaluation
            if question_type == 'descriptive':
                # Simple keyword matching for now
                keywords = set(correct_answer.lower().split())
                student_keywords = set(student_answer.lower().split())
                matching_keywords = keywords.intersection(student_keywords)
                
                score = int((len(matching_keywords) / len(keywords)) * 100) if keywords else 50
                is_correct = score >= 70

            results.append({
                'questionId': answer.get('questionId', ''),
                    'evaluation': f"Found {len(matching_keywords)} relevant keywords",
                'score': score,
                    'isCorrect': is_correct,
                    'confidenceScore': confidence_score
            })

        return jsonify({
            'results': results,
            'message': 'Evaluation completed successfully'
        })

    except Exception as e:
        print(f"Error in evaluate_answers: {str(e)}")
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    initialize_model()
    # Use 0.0.0.0 to allow connections from other containers
    host = os.getenv('FLASK_HOST', '0.0.0.0')
    port = int(os.getenv('FLASK_PORT', 8000))
    debug = os.getenv('FLASK_DEBUG', 'False').lower() == 'true'
    app.run(
        host=host,
        port=port,
        debug=debug
    )